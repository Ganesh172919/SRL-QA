from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Paragraph, Table, TableStyle
from reportlab.pdfgen import canvas

from build_presentation_assets import generate_assets
from presentation.builder import build_presentation_context, build_slide_specs


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

RGB = {
    "navy": RGBColor(0x12, 0x32, 0x4A),
    "teal": RGBColor(0x0F, 0x76, 0x6E),
    "gold": RGBColor(0xC6, 0x8A, 0x2D),
    "ink": RGBColor(0x1F, 0x29, 0x37),
    "slate": RGBColor(0x4B, 0x55, 0x63),
    "cream": RGBColor(0xF7, 0xF3, 0xEB),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}

PDF_COLORS = {
    "navy": colors.HexColor("#12324a"),
    "teal": colors.HexColor("#0f766e"),
    "gold": colors.HexColor("#c68a2d"),
    "ink": colors.HexColor("#1f2937"),
    "slate": colors.HexColor("#4b5563"),
    "cream": colors.HexColor("#f7f3eb"),
    "white": colors.white,
}


def _section_color(section_id: str):
    mapping = {
        "survey": RGB["navy"],
        "llm_integration": RGB["gold"],
        "implementation": RGB["teal"],
        "results_analysis": RGB["gold"],
        "innovation": RGB["slate"],
        "prompt_tuning": RGB["teal"],
        "final_takeaways": RGB["navy"],
        "qna": RGB["navy"],
    }
    return mapping.get(section_id, RGB["navy"])


def _pdf_section_color(section_id: str):
    mapping = {
        "survey": PDF_COLORS["navy"],
        "llm_integration": PDF_COLORS["gold"],
        "implementation": PDF_COLORS["teal"],
        "results_analysis": PDF_COLORS["gold"],
        "innovation": PDF_COLORS["slate"],
        "prompt_tuning": PDF_COLORS["teal"],
        "final_takeaways": PDF_COLORS["navy"],
        "qna": PDF_COLORS["navy"],
    }
    return mapping.get(section_id, PDF_COLORS["navy"])


def _add_background(slide, section_id: str) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGB["cream"]

    top_band = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, Inches(0.7))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = RGB["navy"]
    top_band.line.fill.background()

    accent = slide.shapes.add_shape(1, 0, Inches(0.7), Inches(0.35), Inches(6.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _section_color(section_id)
    accent.line.fill.background()


def _add_title(slide, title: str, section_id: str) -> None:
    textbox = slide.shapes.add_textbox(Inches(0.65), Inches(0.18), Inches(10.8), Inches(0.38))
    frame = textbox.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.name = "Cambria"
    paragraph.font.bold = True
    paragraph.font.size = Pt(24)
    paragraph.font.color.rgb = RGB["white"]
    paragraph.alignment = PP_ALIGN.LEFT

    section_box = slide.shapes.add_textbox(Inches(11.6), Inches(0.18), Inches(1.1), Inches(0.38))
    section_frame = section_box.text_frame
    section_frame.clear()
    section_paragraph = section_frame.paragraphs[0]
    section_paragraph.text = section_id.replace("_", " ").title()
    section_paragraph.font.name = "Aptos"
    section_paragraph.font.bold = True
    section_paragraph.font.size = Pt(10)
    section_paragraph.font.color.rgb = _section_color(section_id)
    section_paragraph.alignment = PP_ALIGN.RIGHT


def _add_footer(slide, text: str, slide_number: int) -> None:
    footer = slide.shapes.add_textbox(Inches(0.7), Inches(7.03), Inches(11.5), Inches(0.25))
    frame = footer.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(9)
    paragraph.font.color.rgb = RGB["slate"]
    paragraph.alignment = PP_ALIGN.LEFT

    number = slide.shapes.add_textbox(Inches(12.2), Inches(7.02), Inches(0.55), Inches(0.26))
    number_frame = number.text_frame
    number_frame.clear()
    number_paragraph = number_frame.paragraphs[0]
    number_paragraph.text = str(slide_number)
    number_paragraph.font.name = "Aptos"
    number_paragraph.font.size = Pt(10)
    number_paragraph.font.bold = True
    number_paragraph.font.color.rgb = RGB["navy"]
    number_paragraph.alignment = PP_ALIGN.RIGHT


def _add_citations(slide, citations: list[str], sources: dict) -> None:
    if not citations:
        return
    entries = []
    for citation_id in citations:
        source = sources.get(citation_id)
        if source:
            entries.append(source)
    if not entries:
        return
    textbox = slide.shapes.add_textbox(Inches(8.0), Inches(0.78), Inches(4.6), Inches(0.32))
    frame = textbox.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(9)
    paragraph.font.color.rgb = RGB["slate"]
    paragraph.text = "Sources: "
    for idx, entry in enumerate(entries):
        run = paragraph.add_run()
        run.text = entry["label"]
        run.font.color.rgb = RGB["teal"]
        run.hyperlink.address = entry["url"]
        if idx < len(entries) - 1:
            sep = paragraph.add_run()
            sep.text = " | "
            sep.font.color.rgb = RGB["slate"]


def _add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float) -> None:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(18 if len(bullets) <= 4 else 16)
        paragraph.font.color.rgb = RGB["ink"]
        paragraph.space_after = Pt(8)
        paragraph.bullet = True


def _add_image(slide, image_path: Path, left: float, top: float, width: float, height: float, caption: str | None) -> None:
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(left), Inches(top + height + 0.05), Inches(width), Inches(0.28))
        caption_frame = caption_box.text_frame
        caption_frame.clear()
        paragraph = caption_frame.paragraphs[0]
        paragraph.text = caption
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(9)
        paragraph.font.color.rgb = RGB["slate"]
        paragraph.alignment = PP_ALIGN.CENTER


def _add_table(slide, headers: list[str], rows: list[list[str]], left: float, top: float, width: float, height: float) -> None:
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    for column_index, header in enumerate(headers):
        cell = table.cell(0, column_index)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGB["navy"]
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(10.5)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGB["white"]
    for row_index, row in enumerate(rows, start=1):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGB["white"]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = RGB["ink"]


def _render_slide(prs: Presentation, spec, figure_paths: dict[str, Path], sources: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, spec.section_id)
    _add_title(slide, spec.title, spec.section_id)
    _add_footer(slide, spec.footer, spec.index)
    _add_citations(slide, spec.citations, sources)

    has_image = spec.image_id is not None
    has_table = bool(spec.table_headers and spec.table_rows)
    if has_image and has_table:
        _add_bullets(slide, spec.bullets, 0.85, 1.05, 4.0, 1.15)
        _add_table(slide, spec.table_headers, spec.table_rows, 0.85, 2.25, 4.2, 3.5)
        _add_image(slide, figure_paths[spec.image_id], 5.45, 1.15, 7.1, 4.8, spec.image_caption)
    elif has_image:
        _add_bullets(slide, spec.bullets, 0.9, 1.15, 4.4, 5.25)
        _add_image(slide, figure_paths[spec.image_id], 5.6, 1.15, 6.55, 4.8, spec.image_caption)
    elif has_table:
        _add_bullets(slide, spec.bullets, 0.9, 1.0, 11.6, 0.9)
        _add_table(slide, spec.table_headers, spec.table_rows, 0.9, 1.95, 11.4, 4.45)
    else:
        _add_bullets(slide, spec.bullets, 1.0, 1.2, 11.3, 5.5)


def _pdf_styles():
    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="SlideTitle",
            parent=styles["Heading1"],
            fontName="Helvetica-Bold",
            fontSize=22,
            leading=26,
            textColor=PDF_COLORS["white"],
        )
    )
    styles.add(
        ParagraphStyle(
            name="SlideBody",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=13,
            leading=17,
            textColor=PDF_COLORS["ink"],
        )
    )
    styles.add(
        ParagraphStyle(
            name="SlideFooter",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=8.5,
            leading=10,
            textColor=PDF_COLORS["slate"],
        )
    )
    return styles


def _draw_paragraph(pdf: canvas.Canvas, text: str, style, x: float, y: float, width: float):
    paragraph = Paragraph(text, style)
    wrapped_width, wrapped_height = paragraph.wrap(width, 1000)
    paragraph.drawOn(pdf, x, y - wrapped_height)
    return wrapped_height


def _draw_pdf_slide(pdf: canvas.Canvas, spec, figure_paths: dict[str, Path], styles, sources: dict) -> None:
    width = 13.333 * inch
    height = 7.5 * inch
    pdf.setFillColor(PDF_COLORS["cream"])
    pdf.rect(0, 0, width, height, fill=1, stroke=0)
    pdf.setFillColor(PDF_COLORS["navy"])
    pdf.rect(0, height - 0.7 * inch, width, 0.7 * inch, fill=1, stroke=0)
    pdf.setFillColor(_pdf_section_color(spec.section_id))
    pdf.rect(0, 0.75 * inch, 0.33 * inch, 6.0 * inch, fill=1, stroke=0)

    _draw_paragraph(pdf, spec.title, styles["SlideTitle"], 0.65 * inch, height - 0.22 * inch, 10.8 * inch)
    pdf.setFillColor(_pdf_section_color(spec.section_id))
    pdf.setFont("Helvetica-Bold", 10)
    pdf.drawRightString(width - 0.55 * inch, height - 0.28 * inch, spec.section_id.replace("_", " ").title())

    has_image = spec.image_id is not None
    has_table = bool(spec.table_headers and spec.table_rows)

    if has_image and has_table:
        x_left = 0.85 * inch
        y_cursor = height - 1.05 * inch
        for bullet in spec.bullets:
            consumed = _draw_paragraph(pdf, f"• {bullet}", styles["SlideBody"], x_left, y_cursor, 4.0 * inch)
            y_cursor -= consumed + 0.08 * inch

        table_data = [spec.table_headers] + spec.table_rows
        table = Table(table_data, colWidths=[4.2 * inch / len(spec.table_headers)] * len(spec.table_headers))
        table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), PDF_COLORS["navy"]),
                    ("TEXTCOLOR", (0, 0), (-1, 0), PDF_COLORS["white"]),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#d0d7de")),
                    ("BACKGROUND", (0, 1), (-1, -1), colors.white),
                    ("FONTSIZE", (0, 0), (-1, -1), 8),
                ]
            )
        )
        tw, th = table.wrapOn(pdf, 4.2 * inch, 3.5 * inch)
        table.drawOn(pdf, x_left, 2.0 * inch)
        image_path = figure_paths[spec.image_id]
        pdf.drawImage(str(image_path), 5.45 * inch, 2.05 * inch, width=7.0 * inch, height=4.5 * inch, preserveAspectRatio=True, mask="auto")
    elif has_image:
        x_left = 0.9 * inch
        y_cursor = height - 1.15 * inch
        for bullet in spec.bullets:
            consumed = _draw_paragraph(pdf, f"• {bullet}", styles["SlideBody"], x_left, y_cursor, 4.2 * inch)
            y_cursor -= consumed + 0.09 * inch
        image_path = figure_paths[spec.image_id]
        pdf.drawImage(str(image_path), 5.55 * inch, 2.0 * inch, width=6.55 * inch, height=4.7 * inch, preserveAspectRatio=True, mask="auto")
    elif has_table:
        x_left = 0.9 * inch
        y_cursor = height - 1.0 * inch
        for bullet in spec.bullets:
            consumed = _draw_paragraph(pdf, f"• {bullet}", styles["SlideBody"], x_left, y_cursor, 11.1 * inch)
            y_cursor -= consumed + 0.08 * inch
        table_data = [spec.table_headers] + spec.table_rows
        table = Table(table_data, colWidths=[11.2 * inch / len(spec.table_headers)] * len(spec.table_headers))
        table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), PDF_COLORS["navy"]),
                    ("TEXTCOLOR", (0, 0), (-1, 0), PDF_COLORS["white"]),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#d0d7de")),
                    ("BACKGROUND", (0, 1), (-1, -1), colors.white),
                    ("FONTSIZE", (0, 0), (-1, -1), 8.2),
                ]
            )
        )
        table.wrapOn(pdf, 11.2 * inch, 4.5 * inch)
        table.drawOn(pdf, 0.9 * inch, 1.6 * inch)
    else:
        x_left = 1.0 * inch
        y_cursor = height - 1.2 * inch
        for bullet in spec.bullets:
            consumed = _draw_paragraph(pdf, f"• {bullet}", styles["SlideBody"], x_left, y_cursor, 11.0 * inch)
            y_cursor -= consumed + 0.11 * inch

    if spec.image_caption:
        _draw_paragraph(pdf, spec.image_caption, styles["SlideFooter"], 5.5 * inch if has_image else 0.9 * inch, 1.55 * inch, 6.7 * inch)

    _draw_paragraph(pdf, spec.footer, styles["SlideFooter"], 0.7 * inch, 0.38 * inch, 6.0 * inch)
    pdf.setFont("Helvetica-Bold", 9)
    pdf.setFillColor(PDF_COLORS["navy"])
    pdf.drawRightString(width - 0.5 * inch, 0.38 * inch, str(spec.index))
    if spec.citations:
        citation_entries = []
        for citation_id in spec.citations:
            source = sources.get(citation_id)
            if source:
                citation_entries.append(f"{source['label']}: {source['url']}")
        if citation_entries:
            pdf.setFont("Helvetica", 8.5)
            pdf.setFillColor(PDF_COLORS["slate"])
            pdf.drawString(0.9 * inch, 0.6 * inch, "Sources: " + " | ".join(citation_entries))


def build_deck() -> dict[str, str]:
    generate_assets()
    ctx = build_presentation_context()
    slide_specs = build_slide_specs(ctx)
    if len(slide_specs) != 40:
        raise RuntimeError(f"Expected 40 slides, found {len(slide_specs)}.")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    for spec in slide_specs:
        _render_slide(prs, spec, ctx.figure_paths, ctx.sources)
    prs.save(ctx.outputs["slide_deck_pptx"])

    pdf = canvas.Canvas(str(ctx.outputs["slide_deck_pdf"]), pagesize=(13.333 * inch, 7.5 * inch))
    styles = _pdf_styles()
    for spec in slide_specs:
        _draw_pdf_slide(pdf, spec, ctx.figure_paths, styles, ctx.sources)
        pdf.showPage()
    pdf.save()

    summary = {
        "slide_count": len(slide_specs),
        "pptx": str(ctx.outputs["slide_deck_pptx"].relative_to(ctx.root)),
        "pdf": str(ctx.outputs["slide_deck_pdf"].relative_to(ctx.root)),
        "titles": [spec.title for spec in slide_specs],
    }
    (ctx.outputs["deck_dir"] / "slide_manifest.json").write_text(json.dumps(summary, indent=2), encoding="utf-8")
    return summary


def main() -> None:
    summary = build_deck()
    print(json.dumps(summary, indent=2))


if __name__ == "__main__":
    main()
