"""Reproducible end-to-end analysis for the local RAISE-SRL-QA project.

Run from the project root:
    python output/code/full_analysis_pipeline.py

All values in generated tables/plots/docs are computed from local files or
local model runs. Use --skip-slow-models only for a partial fast rerun.
"""
from __future__ import annotations

import argparse
import ast
import hashlib
import json
import math
import os
import re
import shutil
import sys
import time
from collections import Counter
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from datasets import Dataset
from pptx import Presentation
from pptx.util import Inches, Pt
from safetensors import safe_open
from sklearn.metrics import confusion_matrix, precision_recall_curve

TOKEN_RE = re.compile(r"[A-Za-z0-9]+(?:[-'][A-Za-z0-9]+)?")
ROLE_RE = re.compile(r"\bARGM-[A-Za-z0-9]+|\bARG-[A-Za-z0-9]+|\bARG[0-9A-Z]+(?:-[A-Za-z0-9]+)?")
STOPWORDS = {
    "the", "and", "for", "that", "with", "was", "were", "are", "this", "from",
    "what", "who", "when", "where", "why", "how", "something", "someone",
    "input", "output", "sentence", "question", "definition", "example",
    "positive", "negative", "complete", "following",
}


@dataclass(frozen=True)
class Paths:
    root: Path
    out: Path
    plots: Path
    tables: Path
    docs: Path
    ppt: Path
    code: Path
    baseline: Path


def paths(project_root: Path | None = None) -> Paths:
    root = (project_root or Path(__file__).resolve().parents[2]).resolve()
    out = root / "output"
    p = Paths(root, out, out / "plots", out / "tables", out / "docs", out / "ppt", out / "code", root.parent / "srl_qa_project")
    for d in (p.out, p.plots, p.tables, p.docs, p.ppt, p.code):
        d.mkdir(parents=True, exist_ok=True)
    return p


def clean_generated_outputs(p: Paths) -> None:
    for folder in (p.plots, p.tables, p.docs, p.ppt):
        folder.mkdir(parents=True, exist_ok=True)
        for child in folder.iterdir():
            if child.is_dir():
                shutil.rmtree(child)
            else:
                child.unlink()


def read_json(path: Path) -> Any:
    for enc in ("utf-8", "utf-16", "utf-8-sig"):
        try:
            return json.loads(path.read_text(encoding=enc))
        except (UnicodeError, json.JSONDecodeError):
            pass
    raise ValueError(f"Could not parse JSON: {path}")


def write_json(path: Path, obj: Any) -> None:
    path.write_text(json.dumps(obj, indent=2, sort_keys=True, default=str), encoding="utf-8")


def tokens(text: Any) -> list[str]:
    return TOKEN_RE.findall(str(text).lower())


def norm(text: Any) -> str:
    return " ".join(tokens(text))


def exact_match(pred: Any, gold: Any) -> float:
    return float(norm(pred) == norm(gold))


def token_prf(pred: Any, gold: Any) -> tuple[float, float, float]:
    pt, gt = tokens(pred), tokens(gold)
    if not pt and not gt:
        return 1.0, 1.0, 1.0
    if not pt or not gt:
        return 0.0, 0.0, 0.0
    overlap = sum((Counter(pt) & Counter(gt)).values())
    if overlap == 0:
        return 0.0, 0.0, 0.0
    prec, rec = overlap / len(pt), overlap / len(gt)
    return prec, rec, 2 * prec * rec / (prec + rec)


def bleu(pred: Any, gold: Any, max_n: int = 4) -> float:
    pt, gt = tokens(pred), tokens(gold)
    if not pt or not gt:
        return 0.0
    scores = []
    for n in range(1, max_n + 1):
        png = Counter(tuple(pt[i:i + n]) for i in range(max(len(pt) - n + 1, 0)))
        gng = Counter(tuple(gt[i:i + n]) for i in range(max(len(gt) - n + 1, 0)))
        scores.append(1.0 if not png else (sum((png & gng).values()) + 1) / (sum(png.values()) + 1))
    bp = 1.0 if len(pt) > len(gt) else math.exp(1 - len(gt) / max(len(pt), 1))
    return float(bp * math.exp(sum(math.log(s) for s in scores) / max_n))


def compact(text: Any) -> str:
    return re.sub(r"\s+", " ", str(text)).strip()


def first_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value.strip()
    if isinstance(value, list):
        return next((first_text(v) for v in value if first_text(v)), "")
    if isinstance(value, dict):
        return next((first_text(value[k]) for k in ("text", "answer", "output") if k in value and first_text(value[k])), "")
    return str(value).strip()


def question_text(value: Any) -> str:
    if isinstance(value, list):
        return compact(" ".join(str(x) for x in value if str(x).strip() and str(x).strip() != "_").replace(" ?", "?"))
    return compact(value)


def role_norm(role: str) -> str:
    return role.strip().upper().replace("ARG-", "ARG")


def roles(text: Any) -> list[str]:
    return [role_norm(r) for r in ROLE_RE.findall(str(text)) if role_norm(r)]


def qtype(text: Any) -> str:
    q = " ".join(tokens(text))
    if not q:
        return "UNKNOWN"
    if re.search(r"\b(to whom|whom)\b", q):
        return "TO-WHOM"
    if re.search(r"\bhow much\b", q):
        return "HOW MUCH"
    for w in ("where", "when", "why", "how", "who", "what"):
        if re.search(rf"\b{w}\b", q):
            return w.upper()
    return "OTHER"


def inferred_role(qt: str) -> str:
    return {
        "WHO": "ARG0", "WHAT": "ARG1", "TO-WHOM": "ARG2",
        "WHERE": "ARGM-LOC", "WHEN": "ARGM-TMP", "WHY": "ARGM-CAU", "HOW": "ARGM-MNR",
    }.get(qt, "")


def is_missing(v: Any) -> bool:
    if v is None:
        return True
    if isinstance(v, str):
        return v.strip() == ""
    if isinstance(v, list):
        return len(v) == 0 or all(is_missing(x) for x in v)
    if isinstance(v, dict):
        return len(v) == 0
    return False


def category(path: Path, root: Path) -> str:
    rel = str(path.relative_to(root)).lower().replace("\\", "/")
    ext = path.suffix.lower()
    if ".hf_cache/datasets/" in rel or rel.startswith("data/") or ext in {".arrow", ".csv"}:
        return "dataset"
    if ".hf_cache/models/" in rel or ext in {".safetensors", ".bin", ".model"}:
        return "model"
    if ext == ".py":
        return "script"
    if ext in {".md", ".txt", ".rst"}:
        return "documentation"
    if ext in {".png", ".jpg", ".jpeg", ".svg"}:
        return "plot_or_image"
    if ext in {".pptx", ".ppt"}:
        return "presentation"
    if "results/" in rel or ext in {".json", ".jsonl"}:
        return "result_or_config"
    if ext == ".log":
        return "log"
    return "other"


def inventory(p: Paths) -> pd.DataFrame:
    rows = []
    for f in p.root.rglob("*"):
        if not f.is_file():
            continue
        rel = str(f.relative_to(p.root)).replace("\\", "/")
        if rel.startswith("output/") and not rel.startswith("output/code/"):
            continue
        st = f.stat()
        rows.append({
            "relative_path": str(f.relative_to(p.root)),
            "absolute_path": str(f),
            "extension": f.suffix.lower() or "[none]",
            "size_bytes": st.st_size,
            "modified_time": datetime.fromtimestamp(st.st_mtime).isoformat(timespec="seconds"),
            "category": category(f, p.root),
        })
    return pd.DataFrame(rows).sort_values(["category", "relative_path"]).reset_index(drop=True)


def parse_final_input(text: str) -> tuple[str, str, str]:
    qs = re.findall(r"Input:\s*Sentence:\s*(.*?)\s*\n\s*Question:\s*(.*?)\s*\nOutput:", text, flags=re.I | re.S)
    if qs:
        s, q = qs[-1]
        return compact(s), compact(q), ""
    vs = re.findall(r"Input:\s*Sentence:\s*(.*?)\s*\n\s*Verb:\s*(.*?)\s*\nOutput:", text, flags=re.I | re.S)
    if vs:
        s, v = vs[-1]
        return compact(s), "", compact(v)
    return "", "", ""


def arrow_split(path: Path, info: dict[str, Any]) -> str:
    name = path.stem.lower()
    for split in sorted(info.get("splits", {}).keys(), key=len, reverse=True):
        if split.lower() in name:
            return split
    for split in ("validation", "valid", "train", "test"):
        if split in name:
            return split
    return "unknown"


def normalize_rows(ds: Dataset, dname: str, split: str, source: str) -> list[dict[str, Any]]:
    out = []
    for i, row in enumerate(ds):
        context = question = answer = predicate = ""
        pred_idx: Any = ""
        role_labels: list[str] = []
        item_count = 0
        if dname == "qa_srl_promptsource":
            context = compact(row.get("sentence", ""))
            question = question_text(row.get("question", ""))
            answer = compact(row.get("rendered_output", "")) or first_text(row.get("answers"))
            predicate = compact(row.get("predicate", ""))
            pred_idx = row.get("predicate_idx", "")
            item_count = len(row.get("answers") or [])
        elif dname == "task1520_qa_srl_answer_generation":
            context, question, predicate = parse_final_input(row.get("input", ""))
            answer = first_text(row.get("output"))
            item_count = len(row.get("output") or [])
        elif dname == "task1519_qa_srl_question_generation":
            context, _, predicate = parse_final_input(row.get("input", ""))
            outputs = row.get("output") or []
            question = " | ".join(compact(x) for x in outputs)
            item_count = len(outputs)
        elif dname == "propbank_srl_seq2seq":
            context = compact(row.get("prompt", ""))
            answer = compact(row.get("response", ""))
            m = re.search(r"SRL for \[(.*?)\]", context, flags=re.I)
            predicate = compact(m.group(1)) if m else ""
            role_labels = roles(answer)
        else:
            context = compact(row.get("sentence", row.get("input", row.get("prompt", ""))))
            question = compact(row.get("question", ""))
            answer = compact(row.get("answer", row.get("output", row.get("response", ""))))
            predicate = compact(row.get("predicate", row.get("verb", "")))
            role_labels = roles(answer)
        qt = qtype(question)
        text_all = " ".join(x for x in (context, question, answer) if x)
        out.append({
            "dataset_name": dname, "split": split, "row_index": i, "source_path": source,
            "context": context, "question": question, "answer": answer, "predicate": predicate,
            "predicate_idx": pred_idx, "question_type": qt,
            "actual_role_labels": "|".join(role_labels), "inferred_role_from_question": inferred_role(qt),
            "context_token_len": len(tokens(context)), "question_token_len": len(tokens(question)),
            "answer_token_len": len(tokens(answer)), "total_token_len": len(tokens(text_all)),
            "answer_or_output_count": item_count,
        })
    return out


def normalize_challenge(records: list[dict[str, Any]], source: str) -> list[dict[str, Any]]:
    out = []
    for i, row in enumerate(records):
        context, question, answer = compact(row.get("context", "")), compact(row.get("question", "")), compact(row.get("expected_answer", ""))
        qt = compact(row.get("question_type", "")) or qtype(question)
        out.append({
            "dataset_name": "challenge_suite_v2", "split": "seed", "row_index": i, "source_path": source,
            "context": context, "question": question, "answer": answer, "predicate": "", "predicate_idx": "",
            "question_type": qt, "actual_role_labels": compact(row.get("target_role", "")),
            "inferred_role_from_question": inferred_role(qt), "context_token_len": len(tokens(context)),
            "question_token_len": len(tokens(question)), "answer_token_len": len(tokens(answer)),
            "total_token_len": len(tokens(context + " " + question + " " + answer)), "answer_or_output_count": 1,
        })
    return out


def load_datasets(p: Paths) -> tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame, pd.DataFrame]:
    overview, schema, missing, normalized = [], [], [], []
    for arrow in sorted((p.root / ".hf_cache" / "datasets").rglob("*.arrow")):
        info_path = arrow.parent / "dataset_info.json"
        info = read_json(info_path) if info_path.exists() else {}
        dname = info.get("dataset_name", arrow.parent.name)
        split = arrow_split(arrow, info)
        ds = Dataset.from_file(str(arrow))
        rel = str(arrow.relative_to(p.root))
        overview.append({
            "dataset_name": dname, "split": split, "rows": len(ds), "columns": len(ds.column_names),
            "file": rel, "size_bytes": arrow.stat().st_size, "builder": info.get("builder_name", ""),
            "version": (info.get("version", {}) or {}).get("version_str", ""),
        })
        first = ds[0] if len(ds) else {}
        for col in ds.column_names:
            vals = ds[col]
            miss = sum(1 for v in vals if is_missing(v))
            row = {
                "dataset_name": dname, "split": split, "column": col,
                "feature": str(ds.features.get(col, "")), "observed_type": type(first.get(col)).__name__ if first else "",
                "missing_count": miss, "missing_rate": miss / max(len(ds), 1), "file": rel,
            }
            schema.append(row)
            missing.append({k: row[k] for k in ("dataset_name", "split", "column", "missing_count", "missing_rate")})
        normalized += normalize_rows(ds, dname, split, rel)
    cpath = p.root / "data" / "challenge_suite_v2.json"
    if cpath.exists():
        recs = read_json(cpath)
        overview.append({
            "dataset_name": "challenge_suite_v2", "split": "seed", "rows": len(recs),
            "columns": len(recs[0]) if recs else 0, "file": str(cpath.relative_to(p.root)),
            "size_bytes": cpath.stat().st_size, "builder": "json", "version": "",
        })
        for col in sorted({k for r in recs for k in r}):
            miss = sum(1 for r in recs if is_missing(r.get(col)))
            schema.append({
                "dataset_name": "challenge_suite_v2", "split": "seed", "column": col, "feature": "json",
                "observed_type": type(recs[0].get(col)).__name__ if recs else "", "missing_count": miss,
                "missing_rate": miss / max(len(recs), 1), "file": str(cpath.relative_to(p.root)),
            })
        normalized += normalize_challenge(recs, str(cpath.relative_to(p.root)))
    return pd.DataFrame(overview), pd.DataFrame(schema), pd.DataFrame(missing), pd.DataFrame(normalized)


def load_frame_store(p: Paths) -> tuple[pd.DataFrame, pd.DataFrame]:
    f = p.root / "retrieval" / "frame_store.json"
    if not f.exists():
        return pd.DataFrame(), pd.DataFrame()
    frames = read_json(f).get("frames", [])
    fr, rr = [], []
    for frame in frames:
        fr.append({
            "lemma": frame.get("lemma", ""), "roleset_id": frame.get("roleset_id", ""),
            "name": frame.get("name", ""), "role_count": len(frame.get("roles", [])),
            "example_count": len(frame.get("examples", [])), "source_file": frame.get("source_file", ""),
        })
        for role in frame.get("roles", []):
            rr.append({"lemma": frame.get("lemma", ""), "roleset_id": frame.get("roleset_id", ""),
                       "role": role_norm(role.get("role", "")), "description": role.get("description", "")})
    return pd.DataFrame(fr), pd.DataFrame(rr)


def load_baseline(p: Paths) -> dict[str, Any]:
    b = p.baseline
    paths_ = {
        "metrics": b / "results" / "metrics.json", "stats": b / "results" / "data_statistics.json",
        "demo": b / "results" / "inference_demo.json", "challenge": b / "data" / "challenge_suite.json",
    }
    out: dict[str, Any] = {"baseline_root_exists": b.exists(), "paths": {k: str(v) for k, v in paths_.items()}}
    for k, f in paths_.items():
        out[k] = read_json(f) if f.exists() else ([] if k in {"demo", "challenge"} else {})
    return out


def baseline_tables(base: dict[str, Any]) -> tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame, pd.DataFrame, pd.DataFrame]:
    m, s = base.get("metrics", {}) or {}, base.get("stats", {}) or {}
    qa, srl, tr = m.get("qa_performance", {}), m.get("srl_performance", {}), m.get("training_diagnostics", {})
    rows = []
    def add(metric: str, value: Any, src: str, scope: str) -> None:
        if value is not None:
            rows.append({"metric": metric, "value": value, "source": src, "scope": scope})
    for key in ("exact_match", "token_f1"):
        add("qa_" + key, qa.get(key), "linked baseline metrics.json", "baseline test")
    for key in ("micro_precision", "micro_recall", "micro_f1", "macro_f1", "bio_accuracy"):
        add("srl_" + key, srl.get(key), "linked baseline metrics.json", "baseline test")
    for key in ("best_epoch", "best_validation_f1", "parameter_count"):
        add(key, tr.get(key), "linked baseline metrics.json", "baseline training")
    for key in ("total_propbank_instances", "usable_propbank_instances", "qa_pair_count", "unique_predicates", "unique_rolesets"):
        add(key, s.get(key), "linked baseline data_statistics.json", "baseline corpus")
    perq = pd.DataFrame([{"question_type": k, **v} for k, v in qa.get("per_question_type", {}).items()])
    err = m.get("error_analysis", {})
    tax = pd.DataFrame([{"error_category": k, "count": v} for k, v in (err.get("taxonomy", {}) or {}).items()])
    role_err = pd.DataFrame([{"role": k, **v} for k, v in (err.get("role_error_rates", {}) or {}).items()])
    return pd.DataFrame(rows), perq, pd.DataFrame(tr.get("history", [])), tax, role_err


def freq_tables(normed: pd.DataFrame, frame_roles: pd.DataFrame, base: dict[str, Any]) -> dict[str, pd.DataFrame]:
    qt = normed["question_type"].replace("", "UNKNOWN").fillna("UNKNOWN").value_counts().rename_axis("question_type").reset_index(name="count")
    rc: Counter[str] = Counter()
    for label_string in normed["actual_role_labels"].fillna("").astype(str):
        for r in label_string.split("|"):
            if r:
                rc[role_norm(r)] += 1
    if not frame_roles.empty:
        rc.update(role_norm(r) for r in frame_roles["role"].dropna().astype(str))
    for r, c in ((base.get("stats", {}) or {}).get("argument_type_distribution", {}) or {}).items():
        rc[role_norm(r)] += int(c)
    words: Counter[str] = Counter()
    for col in ("context", "question", "answer"):
        for text in normed[col].dropna().astype(str):
            words.update(t for t in tokens(text) if len(t) > 2 and t not in STOPWORDS)
    pred = normed["predicate"].replace("", np.nan).dropna().astype(str).str.lower().value_counts().rename_axis("predicate").reset_index(name="count")
    inferred = normed["inferred_role_from_question"].replace("", np.nan).dropna().value_counts().rename_axis("inferred_role_from_question").reset_index(name="count")
    return {
        "question_type_distribution": qt,
        "semantic_role_distribution": pd.DataFrame([{"role": r, "count": c} for r, c in rc.most_common()]),
        "inferred_role_distribution": inferred,
        "top_words": pd.DataFrame([{"token": w, "count": c} for w, c in words.most_common(50)]),
        "predicate_frequency": pred,
    }


def evaluate_seed_suite(p: Paths, skip_slow: bool = False) -> tuple[pd.DataFrame, pd.DataFrame]:
    cpath = p.root / "data" / "challenge_suite_v2.json"
    if not cpath.exists():
        return pd.DataFrame(), pd.DataFrame()
    os.environ.setdefault("HF_DATASETS_OFFLINE", "1")
    os.environ.setdefault("TRANSFORMERS_OFFLINE", "1")
    sys.path.insert(0, str(p.root))
    from srlqa.model_hub import MODEL_SPECS, ModelHub

    keys = ["raise_srlqa_fast"] if skip_slow else [spec.key for spec in MODEL_SPECS]
    hub = ModelHub()
    rows = []
    for ex in read_json(cpath):
        for key in keys:
            res = hub.run_one(key, ex["context"], ex["question"], expected_answer=ex.get("expected_answer"))
            pred, gold = res.get("answer", ""), ex.get("expected_answer", "")
            prec, rec, f1 = token_prf(pred, gold)
            rows.append({
                "example_id": ex.get("id", ""), "diagnostic": ex.get("diagnostic", ""),
                "question_type": ex.get("question_type", qtype(ex.get("question", ""))),
                "target_role": ex.get("target_role", ""), "context": ex.get("context", ""),
                "question": ex.get("question", ""), "expected_answer": gold, "model_key": key,
                "model_label": res.get("model_label", key), "ok": bool(res.get("ok", False)),
                "error": res.get("error", ""), "predicted_answer": pred,
                "predicted_role": res.get("role", ""), "confidence": float(res.get("confidence", 0.0)),
                "latency_ms": float(res.get("latency_ms", 0.0)), "exact_match": exact_match(pred, gold),
                "accuracy": exact_match(pred, gold), "token_precision": prec, "token_recall": rec,
                "token_f1": f1, "bleu": bleu(pred, gold),
                "role_accuracy": float(res.get("role", "") == ex.get("target_role", "")),
                "reasoning": res.get("reasoning", ""),
            })
    df = pd.DataFrame(rows)
    if df.empty:
        return df, pd.DataFrame()
    summary = df.groupby(["model_key", "model_label"], dropna=False).agg(
        examples=("example_id", "count"),
        ok_count=("ok", "sum"),
        error_count=("ok", lambda s: int((~s.astype(bool)).sum())),
        accuracy=("accuracy", "mean"),
        exact_match=("exact_match", "mean"),
        token_precision=("token_precision", "mean"),
        token_recall=("token_recall", "mean"),
        token_f1=("token_f1", "mean"),
        bleu=("bleu", "mean"),
        role_accuracy=("role_accuracy", "mean"),
        confidence_mean=("confidence", "mean"),
        latency_ms_mean=("latency_ms", "mean"),
        latency_ms_median=("latency_ms", "median"),
        latency_ms_p95=("latency_ms", lambda s: float(np.percentile(s, 95))),
    ).reset_index()
    return df, summary


def model_artifacts(p: Paths, summary: pd.DataFrame) -> pd.DataFrame:
    rows = []
    for f in sorted((p.root / ".hf_cache" / "models").rglob("*")) if (p.root / ".hf_cache" / "models").exists() else []:
        if not f.is_file() or f.stat().st_size == 0:
            continue
        params = tensors = np.nan
        if f.suffix.lower() == ".safetensors":
            try:
                total = count = 0
                with safe_open(str(f), framework="pt", device="cpu") as h:
                    for k in h.keys():
                        total += int(np.prod(h.get_tensor(k).shape))
                        count += 1
                params, tensors = total, count
            except Exception:
                pass
        rows.append({
            "artifact": str(f.relative_to(p.root)), "kind": f.suffix.lower().lstrip(".") or "file",
            "size_bytes": f.stat().st_size, "size_mb": f.stat().st_size / (1024 * 1024),
            "parameter_count": params, "tensor_count": tensors,
        })
    df = pd.DataFrame(rows)
    if not summary.empty:
        lat = summary[["model_key", "latency_ms_mean", "latency_ms_median", "latency_ms_p95"]].copy()
        lat["artifact"], lat["kind"] = "measured runtime variant", "runtime"
        df = pd.concat([df, lat], ignore_index=True, sort=False)
    return df


def source_complexity(p: Paths) -> pd.DataFrame:
    rows = []
    for f in sorted((p.root / "srlqa").rglob("*.py")):
        text = f.read_text(encoding="utf-8")
        try:
            tree = ast.parse(text)
            funcs = sum(isinstance(n, (ast.FunctionDef, ast.AsyncFunctionDef)) for n in ast.walk(tree))
            classes = sum(isinstance(n, ast.ClassDef) for n in ast.walk(tree))
        except SyntaxError:
            funcs = classes = 0
        rows.append({"module": str(f.relative_to(p.root)), "lines": len(text.splitlines()), "functions": funcs, "classes": classes})
    return pd.DataFrame(rows)


def ablation(summary: pd.DataFrame) -> pd.DataFrame:
    order = [
        ("legacy_baseline", "Legacy checkpoint baseline"),
        ("legacy_hybrid", "Legacy hybrid additions"),
        ("raise_srlqa_fast", "RAISE retrieval + verifier + deterministic rules"),
        ("raise_srlqa_model", "RAISE plus transformer QA candidates"),
    ]
    rows, base_f1, prev_f1 = [], None, None
    for key, label in order:
        sub = summary[summary["model_key"] == key] if not summary.empty else pd.DataFrame()
        if sub.empty:
            continue
        item = sub.iloc[0]
        if base_f1 is None:
            base_f1 = float(item["token_f1"])
        cur = float(item["token_f1"])
        rows.append({
            "variant": key, "innovation_or_model_change": label, "dataset_scope": "challenge_suite_v2 seed suite",
            "accuracy": float(item["accuracy"]), "token_f1": cur, "role_accuracy": float(item["role_accuracy"]),
            "delta_f1_vs_previous_available_variant": np.nan if prev_f1 is None else cur - prev_f1,
            "delta_f1_vs_legacy_baseline": cur - base_f1,
        })
        prev_f1 = cur
    return pd.DataFrame(rows)


def innovation_matrix(p: Paths, summary: pd.DataFrame) -> pd.DataFrame:
    items = [
        ("PropBank frame retrieval", "srlqa/retrieval/propbank_index.py"),
        ("Constrained span decoding", "srlqa/decoding/span_rules.py"),
        ("Role priors", "srlqa/decoding/role_priors.py"),
        ("Evidence span verifier", "srlqa/verification/span_verifier.py"),
        ("Self-correction loop", "srlqa/verification/self_correction.py"),
        ("Hard negative mining", "srlqa/training/hard_negative_mining.py"),
        ("Distillation scaffolding", "srlqa/distillation/teacher_runner.py"),
        ("Calibrated ensemble", "srlqa/ensemble/weighted_voter.py"),
        ("Nominal event templates", "srlqa/nominal/qa_noun_templates.py"),
        ("Proto-role features", "srlqa/proto_roles/proto_role_features.py"),
    ]
    fast = summary[summary["model_key"] == "raise_srlqa_fast"]["token_f1"] if not summary.empty else pd.Series(dtype=float)
    model = summary[summary["model_key"] == "raise_srlqa_model"]["token_f1"] if not summary.empty else pd.Series(dtype=float)
    rows = []
    for name, rel in items:
        f = p.root / rel
        rows.append({
            "innovation": name, "evidence_file": rel, "file_exists": f.exists(),
            "file_size_bytes": f.stat().st_size if f.exists() else 0,
            "seed_suite_fast_f1_context": float(fast.iloc[0]) if len(fast) else np.nan,
            "seed_suite_model_f1_context": float(model.iloc[0]) if len(model) else np.nan,
            "controlled_ablation_log_found": False,
        })
    return pd.DataFrame(rows)


def sha_small(f: Path) -> str:
    if not f.exists() or f.stat().st_size > 10_000_000:
        return ""
    h = hashlib.sha256()
    with f.open("rb") as fp:
        for chunk in iter(lambda: fp.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def source_manifest(inv: pd.DataFrame) -> pd.DataFrame:
    sub = inv[inv["category"].isin(["dataset", "model", "script", "documentation", "result_or_config"])].copy()
    sub = sub.sort_values(["category", "relative_path"]).head(250)
    sub["sha256_if_le_10mb"] = [sha_small(Path(x)) for x in sub["absolute_path"]]
    return sub


def save_tables(p: Paths, tables: dict[str, pd.DataFrame]) -> None:
    for name, df in tables.items():
        df.to_csv(p.tables / f"{name}.csv", index=False, encoding="utf-8")


def savefig(path: Path) -> None:
    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight")
    plt.close()


def barh(df: pd.DataFrame, x: str, y: str, title: str, path: Path, color: str = "#2f6f9f", xlim: tuple[float, float] | None = None) -> None:
    if df.empty:
        return
    plt.figure(figsize=(11, max(4, min(9, 0.35 * len(df) + 2))))
    ax = sns.barplot(data=df, x=x, y=y, color=color)
    ax.set_title(title)
    ax.set_ylabel("")
    if xlim:
        ax.set_xlim(*xlim)
    savefig(path)


def table_png(df: pd.DataFrame, path: Path, title: str, rows: int = 12, cols: int = 7) -> None:
    shown = df.head(rows).iloc[:, :cols].copy() if not df.empty else pd.DataFrame([{"message": "No rows available"}])
    for col in shown.columns:
        shown[col] = shown[col].map(lambda v: "" if isinstance(v, float) and np.isnan(v) else (f"{v:.4f}" if isinstance(v, float) else str(v)[:70]))
    fig, ax = plt.subplots(figsize=(max(8, 1.8 * len(shown.columns)), max(2.5, 0.45 * (len(shown) + 2))))
    ax.axis("off")
    ax.set_title(title, fontsize=14, fontweight="bold", pad=12)
    tbl = ax.table(cellText=shown.values, colLabels=shown.columns, loc="center", cellLoc="left", colLoc="left")
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(8)
    tbl.scale(1, 1.35)
    for (r, _), cell in tbl.get_celld().items():
        if r == 0:
            cell.set_facecolor("#1f4e79")
            cell.set_text_props(color="white", weight="bold")
        else:
            cell.set_facecolor("#f7fbff" if r % 2 == 0 else "#ffffff")
    savefig(path)


def wrap_label(text: str, width: int = 18) -> str:
    words, lines, cur = text.split(), [], []
    for word in words:
        if sum(len(x) for x in cur) + len(cur) + len(word) > width and cur:
            lines.append(" ".join(cur)); cur = [word]
        else:
            cur.append(word)
    if cur:
        lines.append(" ".join(cur))
    return "\n".join(lines)


def architecture(path: Path, title: str, nodes: dict[str, tuple[float, float]], edges: list[tuple[str, str]]) -> None:
    fig, ax = plt.subplots(figsize=(13, 7))
    ax.axis("off")
    ax.set_title(title, fontsize=16, fontweight="bold", pad=14)
    for name, (x, y) in nodes.items():
        ax.add_patch(plt.Rectangle((x - 0.10, y - 0.055), 0.20, 0.11, facecolor="#e8f3ff", edgecolor="#1f4e79", linewidth=1.8))
        ax.text(x, y, wrap_label(name), ha="center", va="center", fontsize=9, weight="bold")
    for a, b in edges:
        ax.annotate("", xy=nodes[b], xytext=nodes[a], arrowprops=dict(arrowstyle="->", color="#444", lw=1.5, shrinkA=28, shrinkB=28))
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    savefig(path)


def make_architecture_plots(p: Paths) -> dict[str, Path]:
    out: dict[str, Path] = {}
    base = {
        "Context Sentence": (0.10, 0.65), "Question": (0.10, 0.35), "Predicate/Role Parsing": (0.32, 0.50),
        "SRL/QA Model": (0.55, 0.50), "Span Decoder": (0.75, 0.50), "Answer Span + Role": (0.92, 0.50),
    }
    edges = [("Context Sentence", "Predicate/Role Parsing"), ("Question", "Predicate/Role Parsing"), ("Predicate/Role Parsing", "SRL/QA Model"), ("SRL/QA Model", "Span Decoder"), ("Span Decoder", "Answer Span + Role")]
    out["architecture_base_srlqa"] = p.plots / "architecture_base_srlqa.png"; architecture(out["architecture_base_srlqa"], "Base SRL-QA System Architecture", base, edges)
    raise_nodes = {
        "Context + Question": (0.10, 0.50), "Question Type to Role": (0.28, 0.70), "Predicate Inference": (0.28, 0.30),
        "PropBank Retrieval": (0.47, 0.70), "Transformer/Heuristic Candidates": (0.47, 0.30), "Constrained Span Rules": (0.65, 0.30),
        "Evidence Verifier": (0.78, 0.50), "Self-Correction Loop": (0.65, 0.70), "Final Extractive Answer": (0.92, 0.50),
    }
    raise_edges = [("Context + Question", "Question Type to Role"), ("Context + Question", "Predicate Inference"), ("Predicate Inference", "PropBank Retrieval"), ("Question Type to Role", "PropBank Retrieval"), ("Context + Question", "Transformer/Heuristic Candidates"), ("Transformer/Heuristic Candidates", "Constrained Span Rules"), ("Constrained Span Rules", "Evidence Verifier"), ("PropBank Retrieval", "Evidence Verifier"), ("Evidence Verifier", "Self-Correction Loop"), ("Self-Correction Loop", "Evidence Verifier"), ("Evidence Verifier", "Final Extractive Answer")]
    out["architecture_raise_srlqa"] = p.plots / "architecture_raise_srlqa.png"; architecture(out["architecture_raise_srlqa"], "Improved RAISE-SRL-QA Architecture", raise_nodes, raise_edges)
    flow = {"HF Arrow Datasets": (0.10, 0.70), "Local Challenge JSON": (0.10, 0.30), "Normalization": (0.30, 0.50), "EDA Tables": (0.50, 0.70), "Evaluation Runs": (0.50, 0.30), "Plots": (0.70, 0.70), "Docs": (0.88, 0.70), "PPT": (0.88, 0.30)}
    flow_edges = [("HF Arrow Datasets", "Normalization"), ("Local Challenge JSON", "Normalization"), ("Normalization", "EDA Tables"), ("Normalization", "Evaluation Runs"), ("EDA Tables", "Plots"), ("Evaluation Runs", "Plots"), ("Plots", "Docs"), ("Plots", "PPT"), ("Evaluation Runs", "PPT")]
    out["data_flow_diagram"] = p.plots / "data_flow_diagram.png"; architecture(out["data_flow_diagram"], "Reproducible Analysis Data Flow", flow, flow_edges)
    pipe = {"Input Text": (0.09, 0.50), "Candidate Generator": (0.27, 0.50), "Teacher QA Optional": (0.43, 0.72), "Rule Candidates": (0.43, 0.28), "Candidate Dedup": (0.60, 0.50), "Verifier Scoring": (0.76, 0.50), "Best Span": (0.91, 0.50)}
    pipe_edges = [("Input Text", "Candidate Generator"), ("Candidate Generator", "Teacher QA Optional"), ("Candidate Generator", "Rule Candidates"), ("Teacher QA Optional", "Candidate Dedup"), ("Rule Candidates", "Candidate Dedup"), ("Candidate Dedup", "Verifier Scoring"), ("Verifier Scoring", "Best Span")]
    out["model_pipeline_diagram"] = p.plots / "model_pipeline_diagram.png"; architecture(out["model_pipeline_diagram"], "Model Pipeline Diagram", pipe, pipe_edges)
    return out


def make_plots(p: Paths, tables: dict[str, pd.DataFrame]) -> dict[str, Path]:
    sns.set_theme(style="whitegrid")
    out = make_architecture_plots(p)
    inv, overview, schema, normed = tables["file_inventory"], tables["dataset_overview"], tables["column_schema"], tables["normalized_dataset_records"]
    df = inv.groupby("category", as_index=False).agg(files=("relative_path", "count")).sort_values("files", ascending=True)
    out["file_type_distribution"] = p.plots / "file_type_distribution.png"; barh(df, "files", "category", "Project File Categories", out["file_type_distribution"])
    d = overview.assign(dataset_split=overview["dataset_name"] + " / " + overview["split"]).sort_values("rows", ascending=True)
    out["dataset_rows_by_split"] = p.plots / "dataset_rows_by_split.png"; barh(d, "rows", "dataset_split", "Dataset Rows by Source and Split", out["dataset_rows_by_split"], "#277f5f")
    m = schema.assign(field=schema["dataset_name"] + "/" + schema["split"] + "/" + schema["column"]).sort_values("missing_count", ascending=False).head(25).sort_values("missing_count", ascending=True)
    out["missing_values_top"] = p.plots / "missing_values_top.png"; barh(m, "missing_count", "field", "Top Missing-Value Counts by Dataset Field", out["missing_values_top"], "#b45f06")
    lens = normed[["context_token_len", "question_token_len", "answer_token_len", "total_token_len"]].melt(var_name="length_type", value_name="tokens")
    plt.figure(figsize=(11, 6)); ax = sns.histplot(data=lens, x="tokens", hue="length_type", bins=60, element="step", common_norm=False); ax.set_title("Token Length Distributions"); ax.set_xlim(0, min(500, max(1, lens["tokens"].quantile(0.99)))); out["token_length_histogram"] = p.plots / "token_length_histogram.png"; savefig(out["token_length_histogram"])
    box = normed[["dataset_name", "context_token_len", "question_token_len", "answer_token_len"]].melt(id_vars="dataset_name", var_name="length_type", value_name="tokens")
    plt.figure(figsize=(13, 7)); ax = sns.boxplot(data=box, x="tokens", y="dataset_name", hue="length_type", showfliers=False); ax.set_title("Token Length Box Plots by Dataset"); out["token_length_boxplot"] = p.plots / "token_length_boxplot.png"; savefig(out["token_length_boxplot"])
    for key, y, title in [("question_type_distribution", "question_type", "Question Type Distribution"), ("semantic_role_distribution", "role", "Semantic Role Label Frequency"), ("top_words", "token", "Top Word Frequency"), ("predicate_frequency", "predicate", "Top Predicate Frequency")]:
        data = tables[key].head(30).sort_values("count", ascending=True)
        name = {"top_words": "word_frequency_top30", "predicate_frequency": "predicate_frequency_top30"}.get(key, key)
        out[name] = p.plots / f"{name}.png"; barh(data, "count", y, title, out[name], "#2a7f9e")
    plt.figure(figsize=(9, 5)); ax = sns.countplot(data=normed, x="answer_or_output_count", color="#2f6f9f"); ax.set_title("Answer or Output Count Distribution"); out["answer_count_distribution"] = p.plots / "answer_count_distribution.png"; savefig(out["answer_count_distribution"])
    corr = normed[["context_token_len", "question_token_len", "answer_token_len", "total_token_len", "answer_or_output_count"]].corr(numeric_only=True)
    plt.figure(figsize=(7, 6)); ax = sns.heatmap(corr, annot=True, fmt=".2f", cmap="YlGnBu", vmin=-1, vmax=1); ax.set_title("Numeric Feature Correlation Heatmap"); out["correlation_heatmap"] = p.plots / "correlation_heatmap.png"; savefig(out["correlation_heatmap"])
    summ, ev = tables["model_evaluation_summary"], tables["model_evaluation_records"]
    if not summ.empty:
        for metric, fname, title in [("accuracy", "model_accuracy_comparison", "Accuracy / Exact Match by Model"), ("token_f1", "model_f1_comparison", "Token F1 by Model")]:
            data = summ.sort_values(metric, ascending=True); out[fname] = p.plots / f"{fname}.png"; barh(data, metric, "model_label", title, out[fname], "#277f5f", (0, max(1.0, data[metric].max() * 1.05)))
        met = summ[["model_label", "token_precision", "token_recall", "token_f1"]].melt(id_vars="model_label", var_name="metric", value_name="value")
        plt.figure(figsize=(11, 6)); ax = sns.barplot(data=met, x="value", y="model_label", hue="metric"); ax.set_title("Precision, Recall, and F1 by Model"); ax.set_xlim(0, 1.05); out["precision_recall_metric_bars"] = p.plots / "precision_recall_metric_bars.png"; savefig(out["precision_recall_metric_bars"])
        out["latency_comparison"] = p.plots / "latency_comparison.png"; barh(summ.sort_values("latency_ms_mean"), "latency_ms_mean", "model_label", "Mean Inference Latency", out["latency_comparison"], "#b45f06")
    if not ev.empty:
        plt.figure(figsize=(8, 6))
        for label, g in ev.groupby("model_label"):
            y = g["exact_match"].astype(int).to_numpy(); scores = g["confidence"].astype(float).to_numpy()
            if len(set(y.tolist())) >= 2:
                pr, rc, _ = precision_recall_curve(y, scores); plt.plot(rc, pr, marker="o", label=label)
            else:
                plt.scatter([float(y.mean())], [float(y.mean())], s=80, label=f"{label} (single class)")
        plt.title("Precision-Recall Curves from Exact-Match Correctness"); plt.xlabel("Recall"); plt.ylabel("Precision"); plt.xlim(0, 1.05); plt.ylim(0, 1.05); plt.legend(fontsize=8); out["precision_recall_curve"] = p.plots / "precision_recall_curve.png"; savefig(out["precision_recall_curve"])
        labs = sorted(set(ev["target_role"].astype(str)) | set(ev["predicted_role"].astype(str))); labs = [x for x in labs if x]
        cm = confusion_matrix(ev["target_role"], ev["predicted_role"], labels=labs); plt.figure(figsize=(10, 8)); ax = sns.heatmap(cm, annot=True, fmt="d", cmap="YlGnBu", xticklabels=labs, yticklabels=labs); ax.set_title("Role Confusion Matrix Across Seed-Suite Runs"); out["role_confusion_matrix"] = p.plots / "role_confusion_matrix.png"; savefig(out["role_confusion_matrix"])
        plt.figure(figsize=(9, 6)); ax = sns.scatterplot(data=ev, x="confidence", y="token_f1", hue="model_label", style="ok", s=90); ax.set_title("Confidence vs Token F1"); ax.set_xlim(0, 1.05); ax.set_ylim(0, 1.05); out["confidence_vs_f1"] = p.plots / "confidence_vs_f1.png"; savefig(out["confidence_vs_f1"])
    hist = tables["baseline_training_history"]
    if not hist.empty:
        plt.figure(figsize=(9, 5))
        for col in ("train_loss", "validation_loss", "validation_f1"):
            if col in hist: plt.plot(hist["epoch"], hist[col], marker="o", label=col)
        plt.title("Linked Baseline Training Curves"); plt.xlabel("Epoch"); plt.legend(); out["training_loss_curve"] = p.plots / "training_loss_curve.png"; savefig(out["training_loss_curve"])
    perq = tables["baseline_per_question_type"]
    if not perq.empty:
        data = perq[["question_type", "em", "f1"]].melt(id_vars="question_type", var_name="metric", value_name="value")
        plt.figure(figsize=(9, 5)); ax = sns.barplot(data=data, x="value", y="question_type", hue="metric"); ax.set_title("Linked Baseline QA Metrics by Question Type"); ax.set_xlim(0, 1.05); out["baseline_per_question_type"] = p.plots / "baseline_per_question_type.png"; savefig(out["baseline_per_question_type"])
    for key, y, x, title, color in [("baseline_error_taxonomy", "error_category", "count", "Linked Baseline Error Taxonomy", "#b45f06"), ("baseline_role_error_rates", "role", "rate", "Highest Linked Baseline Role Error Rates", "#b45f06")]:
        data = tables[key]
        if not data.empty and x in data:
            data = data.sort_values(x, ascending=False).head(25).sort_values(x, ascending=True)
            out[key] = p.plots / f"{key}.png"; barh(data, x, y, title, out[key], color)
    if not tables["ablation_study"].empty:
        data = tables["ablation_study"][["variant", "accuracy", "token_f1", "role_accuracy"]].melt(id_vars="variant", var_name="metric", value_name="value")
        plt.figure(figsize=(11, 6)); ax = sns.barplot(data=data, x="value", y="variant", hue="metric"); ax.set_title("Available Variant Comparison (Seed Suite)"); ax.set_xlim(0, 1.05); out["ablation_variant_comparison"] = p.plots / "ablation_variant_comparison.png"; savefig(out["ablation_variant_comparison"])
    ma = tables["model_artifacts"]; data = ma[ma["kind"].isin(["safetensors", "bin", "model"])].sort_values("size_mb", ascending=True)
    if not data.empty:
        out["model_artifact_sizes"] = p.plots / "model_artifact_sizes.png"; barh(data, "size_mb", "artifact", "Model Artifact Sizes", out["model_artifact_sizes"], "#2f6f9f")
    for key, title in [("dataset_overview", "Dataset Overview"), ("model_evaluation_summary", "Model Evaluation Summary"), ("baseline_metrics_summary", "Linked Baseline Metrics"), ("ablation_study", "Available Variant Comparison"), ("complexity_summary", "Complexity Summary"), ("innovation_matrix", "Innovation Evidence Matrix"), ("error_cases", "Seed-Suite Error Cases"), ("source_artifact_manifest", "Artifact Inventory Sample")]:
        out[key + "_table"] = p.plots / f"{key}_table.png"; table_png(tables[key], out[key + "_table"], title)
    return {k: v for k, v in out.items() if v.exists()}


def fmt(v: Any) -> str:
    if isinstance(v, float):
        return "" if np.isnan(v) else f"{v:.4f}"
    s = str(v)
    return s[:80] + "..." if len(s) > 83 else s


def md_table(df: pd.DataFrame, rows: int = 10, cols: int = 6) -> str:
    if df.empty:
        return "_No rows available._"
    d = df.head(rows).iloc[:, :cols].copy()
    header = "| " + " | ".join(map(str, d.columns)) + " |"
    sep = "| " + " | ".join("---" for _ in d.columns) + " |"
    body = ["| " + " | ".join(fmt(r[c]).replace("|", "/") for c in d.columns) + " |" for _, r in d.iterrows()]
    return "\n".join([header, sep] + body)


def docs(p: Paths, tables: dict[str, pd.DataFrame], plots: dict[str, Path], meta: dict[str, Any]) -> dict[str, Path]:
    inv, overview, schema = tables["file_inventory"], tables["dataset_overview"], tables["column_schema"]
    evsum, ev = tables["model_evaluation_summary"], tables["model_evaluation_records"]
    baseline, perq = tables["baseline_metrics_summary"], tables["baseline_per_question_type"]
    taxonomy, roleerr = tables["baseline_error_taxonomy"], tables["baseline_role_error_rates"]
    freq_roles, qdist = tables["semantic_role_distribution"], tables["question_type_distribution"]
    best = evsum.sort_values(["token_f1", "accuracy"], ascending=False).iloc[0].to_dict() if not evsum.empty else {}
    lookup = {r["metric"]: r["value"] for _, r in baseline.iterrows()} if not baseline.empty else {}
    file_counts = inv["category"].value_counts().rename_axis("category").reset_index(name="files")
    total_rows = int(overview["rows"].sum()) if not overview.empty else 0
    report = f"""# Complete End-to-End Analysis Report: RAISE-SRL-QA

Generated on: {meta["generated_at"]}

## Abstract

This report analyzes the local `srlqa` project using only files available on disk. The pipeline traversed the project, loaded Hugging Face Arrow datasets, parsed local JSON artifacts, inspected model files, ran available model families on the seed challenge suite, and regenerated tables, plots, documentation, and a PPT under `output/`.

The run found {len(inv)} source/inventory files and loaded {total_rows} rows across dataset artifacts. On `challenge_suite_v2`, the best local measured variant was `{best.get("model_label", "n/a")}` with accuracy/exact match `{best.get("accuracy", np.nan):.4f}` and token F1 `{best.get("token_f1", np.nan):.4f}`. This is a seed-suite result, not an official benchmark claim.

## 1. Introduction

Semantic Role Labeling Question Answering asks natural-language questions about predicate-argument structure and returns extractive answer spans. The local project implements RAISE-SRL-QA: Retrieval-Augmented, Iteratively Self-correcting, Explainable Semantic Role Labeling Question Answering.

## 2. Literature Survey

The codebase combines QA-SRL style question-answer annotations, PropBank-style semantic role inventories, and transformer machine reading comprehension. No external leaderboard values are reported because the requirement is to avoid assumed numbers; every numeric table below comes from local files or local model execution.

## 3. Project Understanding

Inventory summary:

{md_table(file_counts, 12, 2)}

Identified artifact types include dataset files under `.hf_cache/datasets` and `data/`, model files under `.hf_cache/models`, Python scripts under `srlqa/` and root scripts, documentation in Markdown, plots/images, result/config JSON files, and existing PPT decks.

Pipeline flow: input context and question -> question type and role inference -> predicate inference -> PropBank retrieval -> heuristic and optional transformer QA candidates -> constrained span selection -> verifier scoring -> self-correction -> final answer span and role.

## 4. Dataset EDA

Dataset overview:

{md_table(overview.sort_values(["dataset_name", "split"]), 20, 8)}

Top missing-value fields:

{md_table(schema.sort_values("missing_count", ascending=False), 12, 8)}

Question type distribution:

{md_table(qdist, 12, 2)}

Semantic role distribution is computed from artifacts with explicit or parsed role labels: PropBank seq2seq responses, the frame store, the local challenge suite, and linked baseline role statistics. QA-SRL promptsource rows do not expose gold role labels in the cached schema, so inferred roles are saved separately.

{md_table(freq_roles, 15, 2)}

## 5. Methodology

The analysis uses `datasets.Dataset.from_file` for Arrow files, deterministic token normalization for exact match and token F1, a local smoothed sentence-BLEU implementation for text overlap, `safetensors` metadata for parameter counts, and `python-pptx` for the presentation.

## 6. Experiments

The seed-suite evaluation used {len(ev["example_id"].unique()) if not ev.empty else 0} examples from `data/challenge_suite_v2.json`.

{md_table(evsum, 10, 14)}

Linked baseline metrics are a separate scope because they come from `srl_qa_project/results/metrics.json`, not from the seed suite:

{md_table(baseline, 20, 4)}

## 7. Results

Generated figure references include:

- `output/plots/dataset_rows_by_split.png`
- `output/plots/token_length_histogram.png`
- `output/plots/question_type_distribution.png`
- `output/plots/semantic_role_distribution.png`
- `output/plots/model_accuracy_comparison.png`
- `output/plots/model_f1_comparison.png`
- `output/plots/role_confusion_matrix.png`
- `output/plots/latency_comparison.png`

Linked baseline QA EM: `{lookup.get("qa_exact_match", np.nan):.4f}`. Linked baseline QA token F1: `{lookup.get("qa_token_f1", np.nan):.4f}`. Linked baseline SRL micro F1: `{lookup.get("srl_micro_f1", np.nan):.4f}`. Linked baseline SRL macro F1: `{lookup.get("srl_macro_f1", np.nan):.4f}`.

## 8. Error Analysis

Seed-suite error rows are saved to `output/tables/error_cases.csv`. Linked baseline taxonomy:

{md_table(taxonomy, 12, 2)}

Highest linked baseline role error rates:

{md_table(roleerr.sort_values(["rate", "total"], ascending=[False, False]) if not roleerr.empty else roleerr, 12, 4)}

## 9. Ablation and Innovation Analysis

No controlled ablation log file was found for every individual innovation. The table below is therefore an available-variant comparison on the same local seed suite.

{md_table(tables["ablation_study"], 10, 8)}

Innovation evidence:

{md_table(tables["innovation_matrix"], 12, 7)}

## 10. Complexity Analysis

Complexity values are computed from local artifacts: cached model file sizes, safetensors parameter counts, measured seed-suite latency, and Python source line/function/class counts.

{md_table(tables["complexity_summary"], 20, 10)}

## 11. Comparison with Modern QA Systems

The workspace contains cached DeBERTa artifacts (`microsoft/deberta-v3-base` and `deepset/deberta-v3-base-squad2`). External transformer or LLM benchmark scores are not invented here. Future comparison should run those systems through the same `output/code/evaluation_pipeline.py` metrics.

## 12. Conclusion

The project is a reproducible SRL-QA scaffold with retrieval, verification, self-correction, and model-hub components. The generated outputs support local seed-suite comparison and linked-baseline analysis while preserving claim boundaries.

## 13. Future Work

- Add controlled ablations for retrieval, span rules, verification, self-correction, and transformer candidates.
- Freeze a larger held-out benchmark and save prediction JSONL for every model.
- Add calibration plots and reliability analysis.
- Evaluate modern transformer and LLM baselines with the same local metrics.
"""
    report_path = p.docs / "COMPLETE_ANALYSIS_REPORT.md"
    report_path.write_text(report, encoding="utf-8")
    arch = p.docs / "ARCHITECTURE_DIAGRAMS.md"
    arch.write_text("\n\n".join([
        "# Architecture Diagrams",
        f"![Base]({plots['architecture_base_srlqa']})",
        f"![RAISE]({plots['architecture_raise_srlqa']})",
        f"![Data Flow]({plots['data_flow_diagram']})",
        f"![Model Pipeline]({plots['model_pipeline_diagram']})",
    ]), encoding="utf-8")
    repro = p.docs / "REPRODUCIBILITY.md"
    repro.write_text(f"""# Reproducibility Notes

Generated at: {meta["generated_at"]}

Run:

```powershell
python output/code/full_analysis_pipeline.py
```

Source inventory files: {len(inv)}
Dataset rows loaded: {total_rows}
Evaluation records: {len(ev)}
Generated plots: {len(plots)}
Python version: {meta["python_version"]}
""", encoding="utf-8")
    manifest_doc = p.docs / "SOURCE_ARTIFACT_MANIFEST.md"
    manifest_doc.write_text("# Source Artifact Manifest\n\n" + md_table(tables["source_artifact_manifest"], 30, 8), encoding="utf-8")
    return {"report": report_path, "architecture": arch, "reproducibility": repro, "manifest": manifest_doc}


def slide_title(slide: Any, title: str, subtitle: str = "") -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.75))
    p = box.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(30); p.font.bold = True
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.60), Inches(1.08), Inches(11.8), Inches(0.45))
        q = sub.text_frame.paragraphs[0]; q.text = subtitle; q.font.size = Pt(15)


def slide_bullets(slide: Any, bullets: list[str], top: float = 1.7, size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(12.0), Inches(5.2))
    tf = box.text_frame; tf.clear()
    for i, b in enumerate(bullets):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.text = b; par.font.size = Pt(size)


def slide_img(slide: Any, img: Path, top: float = 1.45) -> None:
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0.75), Inches(top), width=Inches(11.8))
    else:
        slide_bullets(slide, [f"Missing image: {img.name}"], top=top)


def pptx(p: Paths, plots: dict[str, Path], tables: dict[str, pd.DataFrame], meta: dict[str, Any]) -> Path:
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5); blank = prs.slide_layouts[6]
    def text(title: str, bullets: list[str], subtitle: str = "") -> None:
        s = prs.slides.add_slide(blank); slide_title(s, title, subtitle); slide_bullets(s, bullets)
    def image(title: str, key: str, subtitle: str = "") -> None:
        s = prs.slides.add_slide(blank); slide_title(s, title, subtitle); slide_img(s, plots.get(key, p.plots / f"{key}.png"))
    best = tables["model_evaluation_summary"].sort_values("token_f1", ascending=False).iloc[0].to_dict() if not tables["model_evaluation_summary"].empty else {}
    total_rows = int(tables["dataset_overview"]["rows"].sum()) if not tables["dataset_overview"].empty else 0
    text("RAISE-SRL-QA Complete Analysis", [f"Project: {p.root}", f"Generated: {meta['generated_at']}", "All values are computed by output/code/full_analysis_pipeline.py"], "Local reproducible EDA, evaluation, documentation, and PPT")
    text("Problem Statement", ["Answer questions about semantic roles in context sentences.", "Return extractive spans with role labels and confidence.", "Prevent hallucinated answers through extractive verification."])
    text("Survey of Existing Work", ["QA-SRL frames semantic roles as QA pairs.", "PropBank provides predicate rolesets and argument labels.", "Transformer extractive QA provides neural span candidates.", "No external leaderboard scores are invented in this deck."])
    text("Gap Analysis", ["Span models can confuse role boundaries and adjunct types.", "QA-only extraction may ignore predicate-role compatibility.", "Reproducibility requires frozen files, metrics, tables, plots, and claim boundaries."])
    text("Proposed System", ["RAISE-SRL-QA: retrieval-augmented, self-correcting, explainable SRL-QA.", "Combines role inference, PropBank retrieval, candidate generation, verification, and correction.", f"Best local seed-suite token F1: {best.get('token_f1', np.nan):.4f}."])
    for title, key in [
        ("Base SRL-QA Architecture", "architecture_base_srlqa"), ("Improved RAISE Architecture", "architecture_raise_srlqa"),
        ("Analysis Data Flow", "data_flow_diagram"), ("Model Pipeline", "model_pipeline_diagram"),
        ("Project File Inventory", "file_type_distribution"),
    ]: image(title, key)
    text("Dataset Overview", [f"Dataset rows loaded: {total_rows}.", f"Dataset overview rows: {len(tables['dataset_overview'])}.", "Schemas and normalized text statistics are saved in output/tables."])
    for title, key in [
        ("Dataset Rows by Split", "dataset_rows_by_split"), ("Dataset Overview Table", "dataset_overview_table"),
        ("Missing Values", "missing_values_top"), ("Token Length Histogram", "token_length_histogram"),
        ("Token Length Box Plot", "token_length_boxplot"), ("Question Type Balance", "question_type_distribution"),
        ("Semantic Role Distribution", "semantic_role_distribution"), ("Word Frequency", "word_frequency_top30"),
        ("Predicate Frequency", "predicate_frequency_top30"), ("Feature Correlation Heatmap", "correlation_heatmap"),
    ]: image(title, key)
    text("Experimental Setup", ["Seed suite: data/challenge_suite_v2.json.", "Metrics: accuracy/EM, precision, recall, token F1, role accuracy, smoothed sentence BLEU.", "Latency measured inside ModelHub.run_one with time.perf_counter."])
    for title, key in [
        ("Accuracy Comparison", "model_accuracy_comparison"), ("F1 Comparison", "model_f1_comparison"),
        ("Metric Summary Table", "model_evaluation_summary_table"), ("Precision and Recall", "precision_recall_metric_bars"),
        ("Precision-Recall Curve", "precision_recall_curve"), ("Role Confusion Matrix", "role_confusion_matrix"),
        ("Latency Comparison", "latency_comparison"), ("Confidence vs F1", "confidence_vs_f1"),
        ("Linked Baseline Training Curve", "training_loss_curve"), ("Linked Baseline Question Types", "baseline_per_question_type"),
        ("Error Taxonomy", "baseline_error_taxonomy"), ("Ablation / Variant Comparison", "ablation_variant_comparison"),
        ("Innovation Evidence Matrix", "innovation_matrix_table"), ("Model Artifact Sizes", "model_artifact_sizes"),
    ]: image(title, key)
    text("Modern QA Comparison", ["Local transformer candidate model: deepset/deberta-v3-base-squad2.", "Local encoder artifact: microsoft/deberta-v3-base.", "External LLM metrics should be evaluated through the same local pipeline, not guessed."])
    text("Conclusion", ["Generated outputs are reproducible from output/code.", "All numeric values come from local files or local model runs.", "Seed-suite scores are not official benchmark claims."])
    text("Future Work", ["Add controlled ablation logs for each innovation.", "Freeze a larger held-out benchmark.", "Run transformer and LLM baselines through the same metrics.", "Add calibration and reliability diagrams."])
    deck = p.ppt / f"RAISE_SRLQA_COMPLETE_ANALYSIS_{len(prs.slides)}_SLIDES.pptx"
    prs.save(deck)
    return deck


def wrappers(p: Paths) -> None:
    files = {
        "data_loading.py": 'from full_analysis_pipeline import paths, inventory, load_datasets\np=paths(); inventory(p).to_csv(p.tables/"file_inventory.csv", index=False); load_datasets(p)[0].to_csv(p.tables/"dataset_overview.csv", index=False)\n',
        "eda.py": 'from full_analysis_pipeline import main\nmain()\n',
        "visualization.py": 'from full_analysis_pipeline import main\nmain()\n',
        "metric_computation.py": 'from full_analysis_pipeline import exact_match, token_prf, bleu\nprint({"exact_match": exact_match("a","a"), "token_prf": token_prf("to office","office"), "bleu": bleu("to office","office")})\n',
        "evaluation_pipeline.py": 'from full_analysis_pipeline import paths, evaluate_seed_suite\np=paths(); r,s=evaluate_seed_suite(p); r.to_csv(p.tables/"model_evaluation_records.csv", index=False); s.to_csv(p.tables/"model_evaluation_summary.csv", index=False)\n',
    }
    for name, text in files.items():
        (p.code / name).write_text(text, encoding="utf-8")


def run(project_root: Path | None = None, skip_slow: bool = False) -> dict[str, Any]:
    p = paths(project_root); wrappers(p); clean_generated_outputs(p); start = time.perf_counter()
    inv = inventory(p)
    overview, schema, missing, normed = load_datasets(p)
    frame_df, frame_roles = load_frame_store(p)
    base = load_baseline(p)
    base_sum, perq, hist, tax, roleerr = baseline_tables(base)
    freqs = freq_tables(normed, frame_roles, base)
    ev, evsum = evaluate_seed_suite(p, skip_slow)
    abl = ablation(evsum)
    innov = innovation_matrix(p, evsum)
    src = source_complexity(p)
    marts = model_artifacts(p, evsum)
    complexity = pd.concat([marts.assign(complexity_source="model/runtime artifacts"), src.assign(complexity_source="python source")], ignore_index=True, sort=False)
    manifest = source_manifest(inv)
    errors = ev[(ev["exact_match"] < 1.0) | (ev["role_accuracy"] < 1.0) | (~ev["ok"].astype(bool))].copy() if not ev.empty else pd.DataFrame()
    tables = {
        "file_inventory": inv, "dataset_overview": overview, "column_schema": schema, "missing_values": missing,
        "normalized_dataset_records": normed, "frame_store_summary": frame_df, "frame_role_distribution_records": frame_roles,
        "baseline_metrics_summary": base_sum, "baseline_per_question_type": perq, "baseline_training_history": hist,
        "baseline_error_taxonomy": tax, "baseline_role_error_rates": roleerr, "model_evaluation_records": ev,
        "model_evaluation_summary": evsum, "ablation_study": abl, "innovation_matrix": innov, "complexity_summary": complexity,
        "model_artifacts": marts, "python_source_complexity": src, "source_artifact_manifest": manifest, "error_cases": errors,
        **freqs,
    }
    save_tables(p, tables); write_json(p.tables / "baseline_artifacts_loaded.json", base)
    meta = {
        "generated_at": datetime.now().isoformat(timespec="seconds"), "project_root": str(p.root),
        "python_version": sys.version.replace("\n", " "), "skip_slow_models": skip_slow,
        "source_inventory_file_count": len(inv), "dataset_rows_loaded": int(overview["rows"].sum()) if not overview.empty else 0,
        "evaluation_records": len(ev), "runtime_seconds": time.perf_counter() - start,
    }
    plots = make_plots(p, tables)
    made_docs = docs(p, tables, plots, meta)
    deck = pptx(p, plots, tables, meta)
    meta.update({"runtime_seconds": time.perf_counter() - start, "plot_count": len(plots), "docs": {k: str(v) for k, v in made_docs.items()}, "ppt": str(deck)})
    write_json(p.tables / "run_metadata.json", meta)
    return {"paths": p, "tables": tables, "plots": plots, "docs": made_docs, "ppt": deck, "metadata": meta}


def main(argv: list[str] | None = None) -> None:
    ap = argparse.ArgumentParser(description="Generate full SRL-QA analysis outputs.")
    ap.add_argument("--project-root", type=Path, default=None)
    ap.add_argument("--skip-slow-models", action="store_true")
    args = ap.parse_args(argv)
    result = run(args.project_root, args.skip_slow_models)
    p: Paths = result["paths"]; meta = result["metadata"]
    print(json.dumps({
        "output_dir": str(p.out), "plots": str(p.plots), "tables": str(p.tables), "docs": str(p.docs),
        "ppt": str(result["ppt"]), "runtime_seconds": meta["runtime_seconds"],
        "dataset_rows_loaded": meta["dataset_rows_loaded"], "evaluation_records": meta["evaluation_records"],
    }, indent=2))


if __name__ == "__main__":
    main()
