from pathlib import Path
import json

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SRLQA = ROOT / "srlqa"
TEMPLATE = SRLQA / "RAISE_SRLQA_PLOTS_FULL_PRESENTATION.pptx"
OUT = SRLQA / "Question_Answering_Using_Semantic_Roles_Conference_35Slides.pptx"

exact = json.loads((SRLQA / "plots" / "EXACT_VALUES.json").read_text(encoding="utf-8"))
model_summary = pd.read_csv(SRLQA / "output" / "tables" / "model_evaluation_summary.csv")
innovation_df = pd.read_csv(SRLQA / "output" / "tables" / "innovation_matrix.csv")
error_taxonomy = pd.read_csv(SRLQA / "output" / "tables" / "baseline_error_taxonomy.csv")

prs = Presentation(str(TEMPLATE))
for idx in range(len(prs.slides) - 1, -1, -1):
    slide_id = prs.slides._sldIdLst[idx]
    prs.part.drop_rel(slide_id.rId)
    del prs.slides._sldIdLst[idx]

BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height

NAVY = RGBColor(6, 18, 30)
NAVY2 = RGBColor(9, 31, 49)
PANEL = RGBColor(14, 45, 67)
PANEL2 = RGBColor(21, 62, 86)
CYAN = RGBColor(70, 210, 220)
GOLD = RGBColor(249, 184, 82)
GREEN = RGBColor(89, 221, 149)
RED = RGBColor(255, 104, 116)
WHITE = RGBColor(246, 250, 255)
MUTED = RGBColor(178, 196, 212)
LINE = RGBColor(71, 100, 122)
TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"
MONO_FONT = "Consolas"
used_images = []


def add_shape(slide, shape_type, x, y, w, h, fill=None, line=None, line_width=1):
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=16,
    color=WHITE,
    bold=False,
    italic=False,
    font=BODY_FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    frame.vertical_anchor = valign
    for idx, line in enumerate(str(text).split("\n")):
        para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        run = para.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def add_deck_background(slide, slide_no, section=""):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, H, NAVY, None)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.10), H, CYAN, None)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.05), GOLD, None)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(10.7), Inches(-1.2), Inches(4.0), Inches(4.0), NAVY2, None)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(11.6), Inches(4.6), Inches(2.6), Inches(2.6), PANEL, None)
    if section:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(0.18), Inches(2.6), Inches(0.34), PANEL2, CYAN, 0.8)
        add_text(slide, section.upper(), Inches(0.48), Inches(0.23), Inches(2.35), Inches(0.22), size=7.8, color=MUTED, bold=True, font=MONO_FONT)
    add_text(slide, f"{slide_no:02d}/35", Inches(11.82), Inches(6.98), Inches(0.95), Inches(0.18), size=8, color=LINE, align=PP_ALIGN.RIGHT, font=MONO_FONT)


def add_title(slide, title, subtitle=None, y=0.62):
    add_text(slide, title, Inches(0.55), Inches(y), Inches(10.7), Inches(0.58), size=30, color=WHITE, bold=True, font=TITLE_FONT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.58), Inches(y + 0.60), Inches(10.4), Inches(0.48), size=12.5, color=MUTED, font=BODY_FONT)


def add_footer_note(slide, text):
    add_text(slide, text, Inches(0.45), Inches(7.02), Inches(9.7), Inches(0.22), size=6.8, color=LINE, font=BODY_FONT)


def add_slide_base(section, slide_no, title, subtitle=None):
    slide = prs.slides.add_slide(BLANK)
    add_deck_background(slide, slide_no, section)
    add_title(slide, title, subtitle)
    return slide


def add_bullets(slide, bullets, x, y, w, h, size=14, color=MUTED, title=None, accent=CYAN, compact=False):
    if title:
        add_text(slide, title, x, y, w, Inches(0.26), size=size + 1.2, color=WHITE, bold=True, font=TITLE_FONT)
        y = y + Inches(0.36)
        h = h - Inches(0.36)
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    for idx, bullet in enumerate(bullets):
        para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        para.line_spacing = 0.90 if compact else 1.05
        para.space_after = Pt(3 if compact else 6)
        text = str(bullet)
        if ":" in text and len(text.split(":", 1)[0]) < 30:
            label, rest = text.split(":", 1)
            r1 = para.add_run()
            r1.text = "- " + label + ":"
            r1.font.name = BODY_FONT
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = accent
            r2 = para.add_run()
            r2.text = rest
            r2.font.name = BODY_FONT
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            run = para.add_run()
            run.text = "- " + text
            run.font.name = BODY_FONT
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def add_card(slide, x, y, w, h, title, body, accent=CYAN, value=None, source=None):
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, PANEL, LINE, 1)
    card.adjustments[0] = 0.12
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h, accent, None)
    add_text(slide, title, x + Inches(0.18), y + Inches(0.13), w - Inches(0.34), Inches(0.28), size=11.5, color=WHITE, bold=True, font=TITLE_FONT)
    if value is not None:
        add_text(slide, value, x + Inches(0.18), y + Inches(0.42), w - Inches(0.34), Inches(0.42), size=22, color=accent, bold=True, font=TITLE_FONT)
        add_text(slide, body, x + Inches(0.18), y + Inches(0.90), w - Inches(0.34), h - Inches(1.08), size=9.5, color=MUTED)
    else:
        add_text(slide, body, x + Inches(0.18), y + Inches(0.46), w - Inches(0.34), h - Inches(0.66), size=10.5, color=MUTED)
    if source:
        add_text(slide, source, x + Inches(0.18), y + h - Inches(0.22), w - Inches(0.34), Inches(0.16), size=6.5, color=LINE, italic=True)
    return card


def add_big_number(slide, value, label, x, y, w, h, accent=CYAN):
    add_card(slide, x, y, w, h, label, "", accent=accent)
    add_text(slide, value, x + Inches(0.18), y + Inches(0.44), w - Inches(0.36), Inches(0.58), size=26, color=accent, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + Inches(0.18), y + Inches(1.03), w - Inches(0.36), Inches(0.28), size=9, color=MUTED, align=PP_ALIGN.CENTER)


def add_picture_fit(slide, path, x, y, w, h, border=True, caption=None):
    path = Path(path)
    if not path.exists():
        add_card(slide, x, y, w, h, "Missing visual", str(path), accent=RED)
        return None
    used_images.append(str(path))
    if border:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x - Inches(0.03), y - Inches(0.03), w + Inches(0.06), h + Inches(0.06), PANEL, LINE, 0.8)
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(int(w) / iw, int(h) / ih)
    new_w, new_h = int(iw * scale), int(ih * scale)
    left = int(x + (int(w) - new_w) / 2)
    top = int(y + (int(h) - new_h) / 2)
    picture = slide.shapes.add_picture(str(path), left, top, width=new_w, height=new_h)
    if caption:
        add_text(slide, caption, x, y + h + Inches(0.05), w, Inches(0.2), size=7.5, color=LINE, align=PP_ALIGN.CENTER)
    return picture


def add_table(slide, rows, x, y, w, h, col_widths=None, font_size=8.5, header_fill=PANEL2, body_fill=PANEL):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    table = shape.table
    if col_widths:
        total = sum(col_widths)
        for idx, col_width in enumerate(col_widths):
            table.columns[idx].width = int(w * col_width / total)
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if row_idx == 0 else body_fill
            cell.text_frame.clear()
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = str(value)
            run.font.name = BODY_FONT
            run.font.size = Pt(font_size if row_idx else font_size + 0.4)
            run.font.bold = row_idx == 0
            run.font.color.rgb = WHITE if row_idx == 0 else MUTED
            para.line_spacing = 0.90
    return shape


def pct(value, digits=1):
    if value == "" or pd.isna(value):
        return "-"
    return f"{float(value) * 100:.{digits}f}%"


def num(value):
    return f"{int(value):,}"


def add_bar_chart(slide, labels, values, x, y, w, h, title, color=CYAN, max_value=None, value_fmt="pct"):
    add_card(slide, x, y, w, h, title, "", accent=color)
    chart_x = x + Inches(0.28)
    chart_y = y + Inches(0.58)
    chart_w = w - Inches(0.55)
    chart_h = h - Inches(0.82)
    max_value = max_value or max(values) or 1
    gap = Inches(0.07)
    bar_h = int((chart_h - gap * (len(values) - 1)) / len(values))
    for idx, (label, value) in enumerate(zip(labels, values)):
        yy = chart_y + idx * (bar_h + gap)
        add_text(slide, label, chart_x, yy + Inches(0.01), Inches(1.58), bar_h, size=7.8, color=MUTED)
        track_x = chart_x + Inches(1.70)
        track_w = chart_w - Inches(2.40)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, track_x, yy + Inches(0.05), track_w, max(Inches(0.08), bar_h - Inches(0.08)), NAVY2, None)
        fill_w = max(Inches(0.02), int(track_w * (value / max_value)))
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, track_x, yy + Inches(0.05), fill_w, max(Inches(0.08), bar_h - Inches(0.08)), color, None)
        val_text = f"{value * 100:.1f}%" if value_fmt == "pct" else f"{value:,.0f}"
        add_text(slide, val_text, track_x + track_w + Inches(0.08), yy + Inches(0.01), Inches(0.60), bar_h, size=7.6, color=WHITE, align=PP_ALIGN.RIGHT, font=MONO_FONT)


def add_pipeline(slide, items, x, y, w, h, accent=CYAN):
    box_w = int((w - Inches(0.24) * (len(items) - 1)) / len(items))
    for idx, item in enumerate(items):
        xx = x + idx * (box_w + Inches(0.24))
        add_card(slide, xx, y, box_w, h, item[0], item[1], accent=accent)
        if idx < len(items) - 1:
            add_shape(slide, MSO_SHAPE.RIGHT_ARROW, xx + box_w - Inches(0.01), y + h / 2 - Inches(0.12), Inches(0.26), Inches(0.24), accent, None)


dc = exact["dataset_counts"]
bm = exact["baseline_metrics"]
qtypes = dc.get("qa_pairs_per_question_type", {})
roles = dc.get("qa_pairs_per_argument_type", {})
order = ["Classical baseline", "Heuristic reranker", "Transformer QA assist", "Full hybrid"]
combined_tracks = [r for r in exact["benchmark_tracks"] if r.get("split_key") == "combined"]
combined_tracks = sorted(combined_tracks, key=lambda row: order.index(row["track"]))
refs_short = {
    "PropBank": "Palmer et al. 2005, ACL Anthology J05-1004",
    "QA-SRL": "He et al. 2015, ACL Anthology D15-1076",
    "Large QA-SRL": "FitzGerald et al. 2018, ACL Anthology P18-1191",
    "Deep SRL": "He et al. 2017, ACL Anthology P17-1044",
    "LISA": "Strubell et al. 2018, ACL Anthology D18-1548",
    "MRC-SRL": "Wang et al. 2022, ACL Anthology 2022.coling-1.191",
    "BERT": "Devlin et al. 2019, ACL Anthology N19-1423",
    "SQuAD": "Rajpurkar et al. 2016, ACL Anthology D16-1264",
    "LoRA": "Hu et al. 2021, arXiv 2106.09685",
    "QLoRA": "Dettmers et al. 2023, arXiv 2305.14314",
}


def build_deck():
    slide = prs.slides.add_slide(BLANK)
    add_deck_background(slide, 1, "Opening")
    add_text(slide, "Question Answering\nUsing Semantic Roles", Inches(0.72), Inches(0.76), Inches(8.4), Inches(1.35), size=37, color=WHITE, bold=True, font=TITLE_FONT, line_spacing=0.87)
    add_text(slide, "A PropBank-grounded SRL-QA system with retrieval, constrained span reasoning, verifier feedback, and PEFT-ready transformer experiments", Inches(0.78), Inches(2.32), Inches(8.1), Inches(0.64), size=14.5, color=MUTED)
    add_big_number(slide, num(dc["qa_pair_count"]), "QA pairs", Inches(0.75), Inches(4.15), Inches(2.05), Inches(1.48), CYAN)
    add_big_number(slide, num(dc["corpus_overview_usable_instances"]), "usable PropBank instances", Inches(3.02), Inches(4.15), Inches(2.35), Inches(1.48), GOLD)
    add_big_number(slide, pct(bm["qa_token_f1"]), "baseline QA token F1", Inches(5.62), Inches(4.15), Inches(2.35), Inches(1.48), GREEN)
    add_card(slide, Inches(9.15), Inches(1.10), Inches(3.35), Inches(4.85), "Core Claim", "SRL-QA improves answerability by making the model reason over who did what to whom, when, where, how, and why - not just over lexical overlap.", accent=GOLD)
    add_footer_note(slide, "Evidence: local project artifacts in srlqa/ plus cited primary research sources.")

    slide = add_slide_base("Problem", 2, "Thesis: SRL turns QA into predicate-argument reasoning", "Instead of asking only where text overlaps, we ask which semantic role answers the question.")
    add_pipeline(slide, [("Question", "Where was the package delivered?"), ("Target role", "WHERE -> ARGM-LOC"), ("Predicate", "delivered"), ("Candidate spans", "to the office; at noon; package"), ("Answer", "to the office")], Inches(0.65), Inches(2.18), Inches(11.7), Inches(1.62), CYAN)
    add_card(slide, Inches(0.8), Inches(4.34), Inches(3.55), Inches(1.35), "Why this matters", "The answer span must be semantically compatible with a role, not merely nearby in the sentence.", accent=CYAN)
    add_card(slide, Inches(4.65), Inches(4.34), Inches(3.55), Inches(1.35), "Research framing", "We formulate SRL as extractive question answering over PropBank-style role structure.", accent=GOLD)
    add_card(slide, Inches(8.5), Inches(4.34), Inches(3.55), Inches(1.35), "Engineering framing", "Use retrieval, span constraints, and verification as lightweight structure around transformer QA.", accent=GREEN)

    slide = add_slide_base("Problem", 3, "Problem statement: text matching is not enough", "SRL-QA must disambiguate semantic roles, predicate scope, and answer span boundaries.")
    add_bullets(slide, ["Role ambiguity: the same entity can be agent, patient, recipient, source, or location across predicates.", "Predicate ambiguity: a question may target a role of a specific predicate, not any nearby event.", "Span boundary errors: QA models often return too much context or omit prepositions needed for role meaning.", "Class imbalance: ARG0/ARG1 dominate, while ARGM and rare core roles remain hard to learn.", "Evaluation sensitivity: exact match is harsh; token F1 can hide wrong-role answers."], Inches(0.72), Inches(1.82), Inches(5.45), Inches(4.35), size=14)
    add_bar_chart(slide, ["Wrong role", "Span boundary", "Other", "Predicate miss"], [685 / 3451, 866 / 3451, 109 / 3451, 2 / 3451], Inches(6.65), Inches(1.72), Inches(5.55), Inches(4.45), "Baseline error pressure", RED)
    add_footer_note(slide, "Counts from srlqa/output/tables/baseline_error_taxonomy.csv; denominator = listed error-analysis total.")

    slide = add_slide_base("Background", 4, "Semantic Role Labeling in one slide", "SRL identifies the predicate and labels the arguments that participate in the event.")
    add_table(slide, [["Role", "Meaning in PropBank-style SRL", "QA mapping"], ["ARG0", "Proto-agent / doer", "Who did it?"], ["ARG1", "Proto-patient / thing affected", "What was affected?"], ["ARG2-ARG5", "Predicate-specific complements", "To whom? from what? for what?"], ["ARGM-TMP", "Temporal modifier", "When?"], ["ARGM-LOC", "Location modifier", "Where?"], ["ARGM-MNR/CAU/PNC", "Manner, cause, purpose", "How? why? for what purpose?"]], Inches(0.72), Inches(1.62), Inches(7.08), Inches(4.62), [1.4, 3.1, 2.0], font_size=9.2)
    add_card(slide, Inches(8.15), Inches(1.78), Inches(3.75), Inches(1.25), "Predicate first", "Arguments are interpreted relative to a predicate such as deliver.01, sell.01, or acquire.01.", accent=CYAN)
    add_card(slide, Inches(8.15), Inches(3.24), Inches(3.75), Inches(1.25), "QA conversion", "A role label becomes a question intent; an argument span becomes the extractive answer.", accent=GOLD)
    add_card(slide, Inches(8.15), Inches(4.70), Inches(3.75), Inches(1.25), "Key risk", "The same surface phrase can be correct for one predicate and wrong for another.", accent=RED)

    slide = add_slide_base("Dataset", 5, "PropBank through NLTK: dataset representation", "The project uses PropBank as the semantic-role backbone and converts roles into QA-style examples.")
    code = "import nltk\nnltk.download('propbank')\nfrom nltk.corpus import propbank\ninstances = propbank.instances()\nroleset = propbank.roleset('deliver.01')"
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(1.68), Inches(5.3), Inches(2.04), NAVY2, CYAN, 1)
    add_text(slide, code, Inches(0.95), Inches(1.92), Inches(4.9), Inches(1.48), size=13.5, color=GREEN, font=MONO_FONT)
    add_bullets(slide, ["Raw view: predicates are anchored in Penn Treebank-style structures and PropBank frame files.", "Clean view: each predicate instance yields a sentence, predicate lemma/roleset, argument labels, and answer spans.", "QA view: a role label becomes a question template; the gold argument span becomes the answer.", "Tree-pointer caveat: NLTK exposes PropBank pointers, so conversion must normalize tree spans to text spans before transformer training."], Inches(6.35), Inches(1.52), Inches(5.55), Inches(3.72), size=12.8)
    add_card(slide, Inches(0.85), Inches(4.35), Inches(2.5), Inches(1.16), "Total instances", "", value=num(dc["total_propbank_instances"]), accent=CYAN)
    add_card(slide, Inches(3.6), Inches(4.35), Inches(2.5), Inches(1.16), "Usable instances", "", value=num(dc["corpus_overview_usable_instances"]), accent=GOLD)
    add_card(slide, Inches(6.35), Inches(4.35), Inches(2.5), Inches(1.16), "Unique predicates", "", value=num(dc["unique_predicates"]), accent=GREEN)
    add_card(slide, Inches(9.1), Inches(4.35), Inches(2.5), Inches(1.16), "Unique rolesets", "", value=num(dc["unique_rolesets"]), accent=RED)
    add_footer_note(slide, "Source: NLTK PropBank HOWTO plus local EXACT_VALUES.json dataset counts.")

    slide = add_slide_base("Goal", 6, "Project goal and evaluation protocol", "Build a reproducible SRL-QA pipeline, then report metrics with claim boundaries.")
    add_bullets(slide, ["Task formulation: extractive question answering over sentences with PropBank-style predicate-role supervision.", "Primary outputs: answer span, inferred role, confidence, and a short reasoning trace.", "Training evidence: baseline PropQA-Net metrics, LoRA/QLoRA notebooks, and local RAISE seed-suite runs.", "Metrics: exact match, token F1, role accuracy or SRL micro F1, latency, and confidence.", "Claim policy: measured results must cite local files; roadmap targets are visibly separated."], Inches(0.72), Inches(1.60), Inches(5.60), Inches(4.70), size=13.3)
    add_big_number(slide, pct(bm["qa_exact_match"]), "baseline QA exact match", Inches(6.75), Inches(1.72), Inches(2.35), Inches(1.42), CYAN)
    add_big_number(slide, pct(bm["qa_token_f1"]), "baseline QA token F1", Inches(9.45), Inches(1.72), Inches(2.35), Inches(1.42), GOLD)
    add_big_number(slide, pct(bm["srl_micro_f1"]), "baseline SRL micro F1", Inches(6.75), Inches(3.52), Inches(2.35), Inches(1.42), GREEN)
    add_big_number(slide, pct(bm["srl_macro_f1"]), "baseline SRL macro F1", Inches(9.45), Inches(3.52), Inches(2.35), Inches(1.42), RED)
    add_footer_note(slide, "Baseline values trace to srlqa/plots/EXACT_VALUES.json and linked baseline metrics.")

    slide = add_slide_base("Survey", 7, "Literature timeline", "The project sits at the intersection of PropBank SRL, QA-SRL, MRC, and efficient fine-tuning.")
    timeline = [("2005", "PropBank", "Predicate-argument corpus"), ("2015", "QA-SRL", "Roles as natural questions"), ("2017", "Deep SRL", "End-to-end neural SRL"), ("2018", "LISA", "Syntax-aware attention"), ("2019", "BERT/SQuAD", "Transformer span QA"), ("2021-23", "LoRA/QLoRA", "PEFT for large models")]
    for idx, (year, name, detail) in enumerate(timeline):
        xx = Inches(0.88 + idx * 1.93)
        add_shape(slide, MSO_SHAPE.OVAL, xx, Inches(2.15), Inches(0.62), Inches(0.62), CYAN if idx % 2 == 0 else GOLD, None)
        add_text(slide, year, xx - Inches(0.15), Inches(1.70), Inches(0.92), Inches(0.28), size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, name, xx - Inches(0.35), Inches(2.92), Inches(1.35), Inches(0.28), size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, detail, xx - Inches(0.58), Inches(3.30), Inches(1.82), Inches(0.62), size=8.6, color=MUTED, align=PP_ALIGN.CENTER)
        if idx < len(timeline) - 1:
            add_shape(slide, MSO_SHAPE.RECTANGLE, xx + Inches(0.62), Inches(2.43), Inches(1.31), Inches(0.05), LINE, None)
    add_card(slide, Inches(0.86), Inches(4.64), Inches(11.12), Inches(1.15), "Survey takeaway", "Earlier work gives us the pieces: PropBank labels, QA-style annotation, transformer span extraction, and efficient adaptation. Our system goal is to combine them with explicit retrieval, constraints, and verification.", accent=GREEN)
    add_footer_note(slide, "Primary survey sources listed on final slide.")

    slide = add_slide_base("Survey", 8, "Survey: PropBank and supervised SRL", "PropBank gives canonical role inventories, but not a ready-made QA model.")
    add_card(slide, Inches(0.78), Inches(1.58), Inches(3.55), Inches(1.45), "What previous work did", "Annotated predicates and their arguments over treebank structures; enabled supervised SRL systems.", accent=CYAN)
    add_card(slide, Inches(4.68), Inches(1.58), Inches(3.55), Inches(1.45), "How models train", "Classic and neural SRL systems learn role labels from gold predicate-argument annotations.", accent=GOLD)
    add_card(slide, Inches(8.58), Inches(1.58), Inches(3.55), Inches(1.45), "Limitation for QA", "Role labels and tree pointers must be converted into natural questions and extractive answer spans.", accent=RED)
    add_bullets(slide, ["Strength: canonical role names make evaluation and error analysis precise.", "Weakness: rare roles and predicate-specific complements create heavy class imbalance.", "Our use: PropBank supplies the role schema and frame retrieval evidence."], Inches(1.0), Inches(3.78), Inches(10.65), Inches(1.45), size=14)
    add_footer_note(slide, refs_short["PropBank"] + "; NLTK PropBank HOWTO.")

    slide = add_slide_base("Survey", 9, "Survey: QA-SRL and natural-language roles", "QA-SRL reframes semantic roles as question-answer pairs written in natural language.")
    add_card(slide, Inches(0.78), Inches(1.62), Inches(3.55), Inches(1.50), "What previous work did", "Asked crowd workers to annotate roles with questions and answer spans instead of formal role labels.", accent=CYAN)
    add_card(slide, Inches(4.68), Inches(1.62), Inches(3.55), Inches(1.50), "How models train", "Parsers learn to generate or recover QA-SRL questions and answer spans from text.", accent=GOLD)
    add_card(slide, Inches(8.58), Inches(1.62), Inches(3.55), Inches(1.50), "Limitation for us", "Natural questions are flexible, but mapping them back to canonical PropBank roles can be noisy.", accent=RED)
    add_pipeline(slide, [("Role", "ARGM-LOC"), ("Question", "Where did it happen?"), ("Answer", "to the office"), ("Need", "role consistency")], Inches(1.05), Inches(4.05), Inches(10.6), Inches(1.24), CYAN)
    add_footer_note(slide, refs_short["QA-SRL"] + "; " + refs_short["Large QA-SRL"] + ".")

    slide = add_slide_base("Survey", 10, "Survey: large-scale QA-SRL parsing", "Large-scale QA-SRL work helps with coverage, but introduces parser and template dependencies.")
    add_table(slide, [["Axis", "Previous QA-SRL systems", "Implication for this project"], ["Data", "Question-answer annotations for predicates", "Good fit for QA interface"], ["Training", "Question generation and answer-span prediction", "Useful but not always PropBank-role exact"], ["Strength", "Human-readable semantic roles", "Improves interpretability"], ["Limitation", "Surface-form variability and parser errors", "Needs canonical role normalization"]], Inches(0.78), Inches(1.62), Inches(7.9), Inches(4.30), [1.2, 3.0, 3.2], font_size=9.5)
    add_card(slide, Inches(9.05), Inches(1.84), Inches(3.05), Inches(1.18), "Design response", "Keep natural questions, but bind them to PropBank role priors.", accent=GREEN)
    add_card(slide, Inches(9.05), Inches(3.30), Inches(3.05), Inches(1.18), "Evaluation response", "Score answer span and role consistency separately.", accent=GOLD)
    add_footer_note(slide, refs_short["Large QA-SRL"] + ".")

    slide = add_slide_base("Survey", 11, "Survey: deep SRL and syntax-aware attention", "Neural SRL improved role labeling, but most systems are not designed as QA pipelines.")
    add_card(slide, Inches(0.78), Inches(1.54), Inches(3.55), Inches(1.45), "Deep SRL", "Supervised neural sequence labeling learns predicate-argument structure end to end.", accent=CYAN, source=refs_short["Deep SRL"])
    add_card(slide, Inches(4.68), Inches(1.54), Inches(3.55), Inches(1.45), "LISA", "Self-attention is biased with linguistic/syntactic structure for SRL.", accent=GOLD, source=refs_short["LISA"])
    add_card(slide, Inches(8.58), Inches(1.54), Inches(3.55), Inches(1.45), "Project gap", "These systems label roles, but do not naturally answer user questions with explainable spans.", accent=RED)
    add_bullets(slide, ["Useful lesson: role structure improves generalization beyond local token overlap.", "Useful limitation: a sequence tagger alone is not enough for QA-style output and interaction.", "Our adaptation: wrap SRL structure around extractive QA and post-hoc verification."], Inches(0.95), Inches(3.88), Inches(10.95), Inches(1.55), size=14)

    slide = add_slide_base("Survey", 12, "Survey: MRC-style SRL", "MRC-SRL casts each role query as a reading-comprehension question.")
    add_pipeline(slide, [("Sentence", "Predicate context"), ("Query", "Role-specific question"), ("Encoder", "Transformer MRC"), ("Output", "Answer span")], Inches(0.90), Inches(1.82), Inches(11.05), Inches(1.42), GOLD)
    add_card(slide, Inches(1.02), Inches(4.10), Inches(3.25), Inches(1.35), "What it gives us", "A practical span-prediction formulation compatible with transformer QA checkpoints.", accent=CYAN)
    add_card(slide, Inches(4.95), Inches(4.10), Inches(3.25), Inches(1.35), "Training style", "Role questions are generated or templated, then optimized as answer-span extraction.", accent=GOLD)
    add_card(slide, Inches(8.88), Inches(4.10), Inches(3.25), Inches(1.35), "Limitation", "Template quality and role/query mismatch can cause brittle generalization.", accent=RED)
    add_footer_note(slide, refs_short["MRC-SRL"] + ".")

    slide = add_slide_base("Survey", 13, "Survey: BERT and SQuAD-style extractive QA", "Transformer QA is a strong span extractor, but needs semantic constraints for SRL reliability.")
    add_card(slide, Inches(0.80), Inches(1.62), Inches(3.30), Inches(1.50), "BERT", "Bidirectional transformer pretraining supplies contextual token representations.", accent=CYAN, source=refs_short["BERT"])
    add_card(slide, Inches(4.45), Inches(1.62), Inches(3.30), Inches(1.50), "SQuAD", "Extractive QA objective trains start/end span prediction from questions and passages.", accent=GOLD, source=refs_short["SQuAD"])
    add_card(slide, Inches(8.10), Inches(1.62), Inches(3.80), Inches(1.50), "SRL risk", "A plausible span can be syntactically nearby but semantically the wrong role.", accent=RED)
    add_bullets(slide, ["Why use it: pretrained QA models know how to localize answer spans.", "Why adapt it: SRL requires predicate and role compatibility, not only answerability.", "Project strategy: transformer candidates are scored by PropBank retrieval and verifier checks."], Inches(1.0), Inches(4.03), Inches(10.65), Inches(1.4), size=14)

    slide = add_slide_base("Survey", 14, "Survey: LoRA/QLoRA and 2B-model feasibility", "Parameter-efficient fine-tuning makes larger language models realistic in a Colab-style workflow.")
    add_card(slide, Inches(0.80), Inches(1.55), Inches(3.45), Inches(1.50), "LoRA", "Train low-rank adapters while keeping the base model frozen; lowers trainable parameters.", accent=CYAN, source=refs_short["LoRA"])
    add_card(slide, Inches(4.62), Inches(1.55), Inches(3.45), Inches(1.50), "QLoRA", "Fine-tune adapters on quantized base weights; improves memory feasibility for larger models.", accent=GOLD, source=refs_short["QLoRA"])
    add_card(slide, Inches(8.45), Inches(1.55), Inches(3.45), Inches(1.50), "2B path", "Gemma 2 2B IT is configured for Colab GPU experiments with quantized loading.", accent=GREEN, source="Gemma 2 2B IT model card")
    add_bullets(slide, ["Notebook evidence: local smoke tests validate the pipeline; full 2B QLoRA execution requires GPU and model access.", "Engineering benefit: adapters keep the implementation lightweight and reproducible.", "Research caution: PEFT enables experiments; it does not guarantee a target accuracy by itself."], Inches(0.96), Inches(4.02), Inches(10.75), Inches(1.46), size=13.8)

    slide = add_slide_base("Survey", 15, "Survey synthesis: what previous work contributes", "We use prior work as building blocks, not as unsupported leaderboard claims.")
    add_table(slide, [["Work", "Training style", "Strength", "Limitation for SRL-QA", "Our use"], ["PropBank", "Gold predicate-role annotation", "Canonical roles", "Not QA formatted", "Frame schema + retrieval"], ["QA-SRL", "Question-answer annotations", "Natural interface", "Role mapping variability", "QA conversion idea"], ["Deep/LISA SRL", "Supervised role labeling", "Strong role structure", "Not interactive QA", "Role-aware constraints"], ["MRC-SRL", "Role query -> span extraction", "Transformer-compatible", "Template dependence", "Task formulation"], ["LoRA/QLoRA", "Adapter fine-tuning", "Low memory", "Needs careful evaluation", "Colab-scale experiments"]], Inches(0.48), Inches(1.42), Inches(12.10), Inches(4.86), [1.35, 2.2, 1.75, 2.25, 1.9], font_size=7.9)
    add_footer_note(slide, "Sources: Palmer 2005; He 2015; FitzGerald 2018; He 2017; Strubell 2018; Wang 2022; Hu 2021; Dettmers 2023.")

    slide = add_slide_base("Survey Analysis", 16, "Gap analysis for our project goal", "The opportunity is not just a better encoder - it is structure around the encoder.")
    add_bullets(slide, ["Canonical role grounding: QA-SRL is readable, but PropBank roles are better for controlled evaluation.", "Answer-span discipline: transformer QA needs constraints for role-compatible span boundaries.", "Evidence retrieval: frame definitions and role priors can narrow plausible answers before final selection.", "Self-correction: a verifier can reject spans that are extractive but semantically wrong.", "Claim discipline: small seed suites are useful engineering tests, not general benchmark proof."], Inches(0.78), Inches(1.58), Inches(5.65), Inches(4.55), size=13.4)
    add_card(slide, Inches(7.05), Inches(1.82), Inches(4.45), Inches(1.15), "Research gap", "Bridge formal SRL, natural-language QA, and verifiable span extraction.", accent=CYAN)
    add_card(slide, Inches(7.05), Inches(3.24), Inches(4.45), Inches(1.15), "Engineering gap", "Make the pipeline simple enough to run locally, but extensible for PEFT models.", accent=GOLD)
    add_card(slide, Inches(7.05), Inches(4.66), Inches(4.45), Inches(1.15), "Presentation gap", "Report measured values, target bands, and evidence scopes separately.", accent=GREEN)

    slide = add_slide_base("Innovation", 17, "Innovation opportunity: RAISE-SRL-QA", "Retrieval-Augmented, Iteratively Self-correcting, Explainable SRL Question Answering.")
    add_pipeline(slide, [("R", "Retrieve PropBank frames"), ("A", "Align question to role"), ("I", "Infer predicate + candidates"), ("S", "Score with verifier"), ("E", "Explain + self-correct")], Inches(0.72), Inches(1.86), Inches(11.75), Inches(1.55), GREEN)
    add_card(slide, Inches(0.95), Inches(4.20), Inches(3.35), Inches(1.28), "Core novelty", "Use symbolic role evidence to constrain neural or heuristic answer candidates.", accent=CYAN)
    add_card(slide, Inches(4.72), Inches(4.20), Inches(3.35), Inches(1.28), "Practical novelty", "Keep a fast deterministic path plus optional transformer-assisted candidate generation.", accent=GOLD)
    add_card(slide, Inches(8.49), Inches(4.20), Inches(3.35), Inches(1.28), "Evaluation novelty", "Separate span correctness, role correctness, latency, and evidence scope.", accent=GREEN)

    slide = add_slide_base("Implementation", 18, "Existing folder analysis: what is already implemented", "The current project contains a real SRL-QA scaffold plus reports, plots, and model/evaluation artifacts.")
    rows = [["Component", "Evidence file", "Status"]]
    for _, row in innovation_df.head(9).iterrows():
        rows.append([row["innovation"], row["evidence_file"], "Present" if bool(row["file_exists"]) else "Missing"])
    add_table(slide, rows, Inches(0.52), Inches(1.45), Inches(12.0), Inches(4.95), [2.6, 4.6, 1.0], font_size=7.8)
    add_footer_note(slide, "Source: srlqa/output/tables/innovation_matrix.csv.")

    slide = add_slide_base("Implementation", 19, "Data pipeline: PropBank to SRL-QA examples", "The pipeline turns predicate-argument structure into transformer-ready question-answer instances.")
    add_picture_fit(slide, SRLQA / "plots" / "accurate_assets" / "42_data_pipeline_architecture.png", Inches(0.64), Inches(1.45), Inches(6.30), Inches(4.80), caption="Data pipeline architecture")
    add_bullets(slide, ["Load PropBank/NLTK or cached SRL-QA assets.", "Normalize sentences, predicates, rolesets, and argument spans.", "Generate question templates by role and question type.", "Split into train/validation/test or seed-suite evaluation scope.", "Export tables and plots for reproducible reporting."], Inches(7.25), Inches(1.65), Inches(4.75), Inches(4.35), size=13.2)

    slide = add_slide_base("EDA", 20, "Dataset and EDA: counts, roles, and imbalance", "The role distribution is highly skewed; this explains why macro-F1 remains hard.")
    q_pairs = sorted(qtypes.items(), key=lambda pair: pair[1], reverse=True)
    add_bar_chart(slide, [p[0] for p in q_pairs], [p[1] for p in q_pairs], Inches(0.60), Inches(1.44), Inches(5.75), Inches(4.95), "QA pairs by question type", CYAN, max_value=max(qtypes.values()), value_fmt="count")
    role_pairs = sorted([(k, v) for k, v in roles.items() if isinstance(v, int)], key=lambda pair: pair[1], reverse=True)[:8]
    add_bar_chart(slide, [p[0] for p in role_pairs], [p[1] for p in role_pairs], Inches(6.72), Inches(1.44), Inches(5.75), Inches(4.95), "Top argument labels", GOLD, max_value=max([p[1] for p in role_pairs]), value_fmt="count")
    add_footer_note(slide, "Source: srlqa/plots/EXACT_VALUES.json. Counts are generated from local artifacts.")

    slide = add_slide_base("Implementation", 21, "Architecture: baseline SRL-QA pipeline", "The baseline provides the comparison point for improvements and error analysis.")
    add_picture_fit(slide, SRLQA / "output" / "plots" / "architecture_base_srlqa.png", Inches(0.70), Inches(1.50), Inches(6.05), Inches(4.70), caption="Baseline architecture")
    add_bullets(slide, ["Input: sentence/context plus a natural-language question.", "Model: legacy PropQA-Net or equivalent baseline span predictor.", "Output: answer span with QA exact match/token F1 scoring.", "Observed issue: role errors and span-boundary errors dominate the error taxonomy."], Inches(7.02), Inches(1.72), Inches(4.92), Inches(3.92), size=13.2)
    add_card(slide, Inches(7.10), Inches(5.18), Inches(4.58), Inches(0.78), "Baseline token F1", "Measured saved baseline", value=pct(bm["qa_token_f1"]), accent=GOLD)

    slide = add_slide_base("Implementation", 22, "Architecture: RAISE-SRL-QA pipeline", "The proposed system layers retrieval, constrained decoding, and verification around answer candidates.")
    add_picture_fit(slide, SRLQA / "plots" / "accurate_assets" / "43_raise_architecture.png", Inches(0.64), Inches(1.38), Inches(7.15), Inches(4.95), caption="RAISE-SRL-QA architecture")
    add_bullets(slide, ["Question parser maps WH-form to a likely semantic role.", "Predicate detector selects the target event.", "PropBank frame retrieval narrows legal arguments.", "Candidate generator proposes spans from heuristic and transformer paths.", "Verifier accepts, rejects, or corrects the span."], Inches(8.08), Inches(1.62), Inches(4.02), Inches(4.25), size=12.5)

    slide = add_slide_base("Implementation", 23, "Module detail: PropBank frame retrieval", "Frame retrieval supplies role definitions and predicate-specific constraints.")
    add_picture_fit(slide, SRLQA / "plots" / "accurate_assets" / "45_retrieval_architecture.png", Inches(0.70), Inches(1.42), Inches(6.35), Inches(4.85), caption="Retrieval architecture")
    add_card(slide, Inches(7.34), Inches(1.62), Inches(4.25), Inches(1.12), "Input", "Predicate lemma, roleset candidate, and question-derived role intent.", accent=CYAN)
    add_card(slide, Inches(7.34), Inches(3.05), Inches(4.25), Inches(1.12), "Retrieval object", "Frame roles, role descriptions, examples, and compatible argument types.", accent=GOLD)
    add_card(slide, Inches(7.34), Inches(4.48), Inches(4.25), Inches(1.12), "Why it helps", "It filters spans that look answer-like but conflict with the predicate's semantic frame.", accent=GREEN)

    slide = add_slide_base("Implementation", 24, "Module detail: constrained span decoding and role priors", "The decoder uses role-specific span rules before final answer selection.")
    add_picture_fit(slide, SRLQA / "plots" / "accurate_assets" / "46_constrained_decoding_architecture.png", Inches(0.70), Inches(1.45), Inches(6.35), Inches(4.78), caption="Constrained decoding architecture")
    add_bullets(slide, ["WHERE questions prefer prepositional/location spans over full predicate tails.", "WHEN questions prefer temporal spans such as dates, times, and ARGM-TMP candidates.", "WHO questions prefer entity/person spans compatible with ARG0 or recipient-like roles.", "Role priors reduce over-long answer spans and wrong modifier selection."], Inches(7.34), Inches(1.78), Inches(4.45), Inches(3.65), size=12.6)
    add_footer_note(slide, "Evidence: srlqa/decoding/span_rules.py and srlqa/decoding/role_priors.py.")

    slide = add_slide_base("Implementation", 25, "Module detail: verifier and iterative correction", "The verifier turns a one-shot prediction into an inspectable correction loop.")
    add_picture_fit(slide, SRLQA / "plots" / "accurate_assets" / "47_verifier_correction_architecture.png", Inches(0.68), Inches(1.38), Inches(6.55), Inches(4.90), caption="Verifier and correction architecture")
    add_bullets(slide, ["Check extractability: the answer must occur as a span in the context.", "Check role compatibility: candidate span should match the inferred semantic role.", "Check predicate context: the span should belong to the target event, not a distractor event.", "Self-correct by choosing the next best candidate when verification fails."], Inches(7.54), Inches(1.72), Inches(4.32), Inches(4.05), size=12.7)

    slide = add_slide_base("Experiments", 26, "Model experiments: all available local approaches", "The project compares classical, heuristic, transformer-assisted, and hybrid tracks under a local scope.")
    rows = [["Track", "Scope", "EM", "Token F1", "Role acc.", "Latency"]]
    for row in combined_tracks:
        rows.append([row["track"], "60 ex. local", pct(row["exact_match"]), pct(row["token_f1"]), pct(row["role_accuracy"]), f"{row['mean_latency_ms']:.1f} ms"])
    add_table(slide, rows, Inches(0.72), Inches(1.50), Inches(8.05), Inches(3.12), [2.4, 1.3, 1.0, 1.0, 1.0, 1.2], font_size=9.1)
    add_bullets(slide, ["Classical baseline is fastest but weak on role-sensitive questions.", "Heuristic reranker improves role accuracy without transformer latency.", "Transformer assist and full hybrid match heuristic scores here but cost more latency.", "Interpretation: structural constraints mattered more than model size in this local comparison."], Inches(9.05), Inches(1.50), Inches(3.20), Inches(3.65), size=10.7)
    add_picture_fit(slide, SRLQA / "plots" / "accurate_assets" / "21_combined_improvement_over_classical.png", Inches(0.95), Inches(4.86), Inches(6.0), Inches(1.28), border=True)
    add_footer_note(slide, "Source: srlqa/plots/EXACT_VALUES.json combined benchmark tracks.")

    slide = add_slide_base("Experiments", 27, "LoRA/QLoRA notebook experiments", "The notebooks keep implementation simple while exposing an upgrade path to larger models.")
    add_table(slide, [["Notebook", "Local smoke-test path", "Full experiment path", "Measured local note"], ["LoRA/QLoRA encoder", "DistilBERT QA + LoRA", "DeBERTa-style QA fine-tuning", "test token F1 0.4151; EM 0.1250"], ["2B Gemma QLoRA", "tiny GPT-2 smoke test", "google/gemma-2-2b-it 4-bit QLoRA", "local CPU smoke test did not reach target"]], Inches(0.62), Inches(1.48), Inches(12.0), Inches(1.72), [2.0, 2.3, 2.7, 2.7], font_size=8.2)
    add_card(slide, Inches(0.92), Inches(3.92), Inches(3.40), Inches(1.25), "Why LoRA", "Adapter tuning reduces trainable parameters while keeping base weights frozen.", accent=CYAN)
    add_card(slide, Inches(4.82), Inches(3.92), Inches(3.40), Inches(1.25), "Why QLoRA", "4-bit quantized loading is the practical route for 2B-class Colab experiments.", accent=GOLD)
    add_card(slide, Inches(8.72), Inches(3.92), Inches(3.40), Inches(1.25), "Claim boundary", "The 2B path is implemented as a reproducible Colab-ready experiment, not as a locally proven 80% result.", accent=RED)
    add_footer_note(slide, "Sources: local notebooks plus Gemma 2 2B IT model card.")

    slide = add_slide_base("Results", 28, "Main measured results with evidence scope", "Every numeric result on this slide is tied to a local file and an explicit evaluation scope.")
    rows = [["Evidence scope", "System/metric", "EM", "Token F1", "Role/SRL", "Source"], ["Saved baseline test", "Legacy PropQA-Net", pct(bm["qa_exact_match"]), pct(bm["qa_token_f1"]), pct(bm["srl_micro_f1"]), "EXACT_VALUES.json"], ["60-example local", "Classical baseline", pct(combined_tracks[0]["exact_match"]), pct(combined_tracks[0]["token_f1"]), pct(combined_tracks[0]["role_accuracy"]), "benchmark tracks"], ["60-example local", "Full hybrid", pct(combined_tracks[-1]["exact_match"]), pct(combined_tracks[-1]["token_f1"]), pct(combined_tracks[-1]["role_accuracy"]), "benchmark tracks"], ["15-example seed suite", "Legacy hybrid", pct(float(model_summary.loc[model_summary.model_key == "legacy_hybrid", "exact_match"].iloc[0])), pct(float(model_summary.loc[model_summary.model_key == "legacy_hybrid", "token_f1"].iloc[0])), pct(float(model_summary.loc[model_summary.model_key == "legacy_hybrid", "role_accuracy"].iloc[0])), "model_evaluation_summary.csv"], ["15-example seed suite", "RAISE-SRL-QA Fast", "100.0%", "100.0%", "100.0%", "seed-suite only"]]
    add_table(slide, rows, Inches(0.42), Inches(1.32), Inches(12.34), Inches(4.05), [1.75, 2.0, 0.8, 0.9, 0.9, 1.8], font_size=8.0)
    add_card(slide, Inches(0.84), Inches(5.72), Inches(11.30), Inches(0.58), "Interpretation", "The strongest measured RAISE result is a local 15-example seed-suite result. It is useful engineering evidence, but not an official general benchmark claim.", accent=GOLD)

    slide = add_slide_base("Results", 29, "Result plots: quality, latency, and confidence", "The plots show why structural methods can improve correctness while latency depends on model choice.")
    add_picture_fit(slide, SRLQA / "output" / "plots" / "model_accuracy_comparison.png", Inches(0.62), Inches(1.35), Inches(3.0), Inches(2.18), caption="Accuracy")
    add_picture_fit(slide, SRLQA / "output" / "plots" / "model_f1_comparison.png", Inches(3.80), Inches(1.35), Inches(3.0), Inches(2.18), caption="Token F1")
    add_picture_fit(slide, SRLQA / "output" / "plots" / "latency_comparison.png", Inches(6.98), Inches(1.35), Inches(3.0), Inches(2.18), caption="Latency")
    add_picture_fit(slide, SRLQA / "output" / "plots" / "confidence_vs_f1.png", Inches(10.16), Inches(1.35), Inches(2.45), Inches(2.18), caption="Confidence vs F1")
    add_bullets(slide, ["Fast RAISE path is best for interactive presentation settings under the local seed-suite scope.", "Model-backed path is useful for candidate generation but adds significant latency.", "Confidence should be interpreted with calibration checks, not as proof of correctness."], Inches(0.84), Inches(4.40), Inches(11.3), Inches(1.20), size=13.1)
    add_footer_note(slide, "Plots from srlqa/output/plots; local seed-suite scope where applicable.")

    slide = add_slide_base("Analysis", 30, "Error analysis: where the system still fails", "The baseline error profile motivates the RAISE components.")
    add_picture_fit(slide, SRLQA / "output" / "plots" / "role_confusion_matrix.png", Inches(0.64), Inches(1.36), Inches(5.95), Inches(4.85), caption="Role confusion matrix")
    add_bar_chart(slide, list(error_taxonomy["error_category"]), list(error_taxonomy["count"]), Inches(7.00), Inches(1.42), Inches(5.30), Inches(3.70), "Error taxonomy counts", RED, max_value=max(error_taxonomy["count"]), value_fmt="count")
    add_card(slide, Inches(7.08), Inches(5.35), Inches(5.1), Inches(0.72), "Main lesson", "Span-boundary and wrong-role errors are exactly the errors targeted by constrained decoding and verifier correction.", accent=GOLD)
    add_footer_note(slide, "Source: baseline_error_taxonomy.csv and role_confusion_matrix.png.")

    slide = add_slide_base("Innovation", 31, "Innovation matrix: implemented vs future-evaluable", "The implementation contains multiple innovation modules, but not every module has a controlled ablation yet.")
    rows = [["Innovation", "Implemented evidence", "Independent ablation?"]]
    for _, row in innovation_df.head(8).iterrows():
        rows.append([row["innovation"], row["evidence_file"], "No controlled log" if not bool(row["controlled_ablation_log_found"]) else "Yes"])
    add_table(slide, rows, Inches(0.50), Inches(1.34), Inches(12.15), Inches(4.62), [2.6, 4.7, 1.65], font_size=7.7)
    add_footer_note(slide, "Source: innovation_matrix.csv. Seed-suite context is reported separately from controlled ablations.")

    slide = add_slide_base("Comparison", 32, "Comparison with previous work", "The comparison emphasizes method fit and claim boundaries rather than unsupported leaderboard numbers.")
    add_table(slide, [["Capability", "PropBank SRL", "QA-SRL", "Deep/LISA", "MRC-SRL", "This project"], ["Canonical roles", "Yes", "Partial", "Yes", "Yes", "Yes"], ["Natural QA interface", "No", "Yes", "No", "Yes", "Yes"], ["Transformer span QA", "No", "Partial", "No", "Yes", "Optional"], ["Frame retrieval", "No", "No", "No", "Partial", "Yes"], ["Constrained decoding", "No", "No", "Partial", "Partial", "Yes"], ["Verifier/self-correction", "No", "No", "No", "Partial", "Yes"], ["PEFT 2B path", "No", "No", "No", "No", "Yes"]], Inches(0.45), Inches(1.28), Inches(12.25), Inches(4.70), [1.95, 1.2, 1.2, 1.2, 1.2, 1.35], font_size=7.8)
    add_footer_note(slide, "Prior-work columns summarize cited methods qualitatively; project metrics remain local-file based.")

    slide = add_slide_base("Roadmap", 33, "Target roadmap: separated from measured results", "These are target bands for the next experimental cycle, not achieved scores.")
    add_card(slide, Inches(0.80), Inches(1.58), Inches(3.25), Inches(3.75), "Stage 1", "Stabilize PropBank conversion, span-boundary rules, and frozen benchmark scripts.", value="75%", accent=CYAN, source="target token-F1 band")
    add_card(slide, Inches(4.75), Inches(1.58), Inches(3.25), Inches(3.75), "Stage 2", "Add controlled ablations for retrieval, role priors, verifier, and LoRA fine-tuning.", value="80%", accent=GOLD, source="target token-F1 band")
    add_card(slide, Inches(8.70), Inches(1.58), Inches(3.25), Inches(3.75), "Stage 3", "Scale to 2B QLoRA, hard negatives, calibrated ensemble, and external benchmark validation.", value="85%", accent=GREEN, source="target token-F1 band")
    add_footer_note(slide, "Roadmap values are intentionally separated from measured-result slides.")

    slide = add_slide_base("Limitations", 34, "Limitations and threats to validity", "The deck is intentionally conservative about what has been proven.")
    add_bullets(slide, ["Evaluation scope: several strong RAISE scores come from a 15-example handcrafted seed suite.", "Benchmark gap: no external leaderboard claim is made for the local system.", "Ablation gap: innovation modules are implemented, but some lack independent controlled ablations.", "Dataset conversion risk: PropBank tree pointers require careful text-span normalization.", "Model access risk: full Gemma 2B QLoRA requires GPU runtime and accepted Hugging Face model terms.", "Future fix: freeze a benchmark split, add end-to-end notebook execution logs, and report statistical confidence intervals."], Inches(0.78), Inches(1.55), Inches(6.30), Inches(4.72), size=13.2)
    add_card(slide, Inches(7.55), Inches(1.76), Inches(4.05), Inches(1.25), "What we can claim", "A working, explainable SRL-QA implementation with local evidence and reproducible artifacts.", accent=GREEN)
    add_card(slide, Inches(7.55), Inches(3.38), Inches(4.05), Inches(1.25), "What we should not claim", "General state-of-the-art performance or achieved 80% accuracy without a frozen benchmark run.", accent=RED)
    add_card(slide, Inches(7.55), Inches(5.00), Inches(4.05), Inches(1.00), "Conference posture", "Strong system story, honest evaluation, clear next experiments.", accent=GOLD)

    slide = add_slide_base("Conclusion", 35, "Conclusion: SRL-QA as structured answer extraction", "The contribution is a principled bridge between PropBank roles, QA interaction, and verifiable span extraction.")
    add_card(slide, Inches(0.78), Inches(1.48), Inches(3.35), Inches(1.30), "Research message", "Semantic roles make QA more controllable by linking questions to predicate-specific argument structure.", accent=CYAN)
    add_card(slide, Inches(4.52), Inches(1.48), Inches(3.35), Inches(1.30), "Engineering message", "RAISE-SRL-QA adds retrieval, constraints, verification, and optional PEFT models around a simple QA core.", accent=GOLD)
    add_card(slide, Inches(8.26), Inches(1.48), Inches(3.35), Inches(1.30), "Next step", "Run controlled ablations and the full 2B QLoRA experiment on a frozen benchmark split.", accent=GREEN)
    refs = (
        "Primary references: https://aclanthology.org/J05-1004/ | "
        "https://aclanthology.org/D15-1076/ | https://aclanthology.org/P18-1191/ | "
        "https://aclanthology.org/P17-1044/ | https://aclanthology.org/D18-1548/ | "
        "https://aclanthology.org/2022.coling-1.191/ | https://aclanthology.org/N19-1423/ | "
        "https://aclanthology.org/D16-1264/ | https://arxiv.org/abs/2106.09685 | "
        "https://arxiv.org/abs/2305.14314 | https://www.nltk.org/howto/propbank.html | "
        "https://huggingface.co/google/gemma-2-2b-it"
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(3.42), Inches(11.30), Inches(2.22), NAVY2, LINE, 1)
    add_text(slide, "Reference URLs", Inches(1.02), Inches(3.62), Inches(10.82), Inches(0.26), size=12, color=WHITE, bold=True, font=TITLE_FONT)
    add_text(slide, refs, Inches(1.02), Inches(4.00), Inches(10.75), Inches(1.25), size=7.5, color=MUTED, font=BODY_FONT)
    add_footer_note(slide, "All project numbers in the deck are either local measured values or explicitly labeled roadmap targets.")


def validate_deck():
    prs2 = Presentation(str(OUT))
    assert len(prs2.slides) == 35, f"Expected 35 slides, found {len(prs2.slides)}"
    missing_text = []
    for idx, slide in enumerate(prs2.slides, start=1):
        slide_text = "\n".join([getattr(shape, "text", "") for shape in slide.shapes if hasattr(shape, "text")]).strip()
        if not slide_text:
            missing_text.append(idx)
    assert not missing_text, f"Slides without text: {missing_text}"
    for image_path in used_images:
        assert Path(image_path).exists(), f"Missing image used in deck: {image_path}"
    all_text = "\n".join(
        "\n".join([getattr(shape, "text", "") for shape in slide.shapes if hasattr(shape, "text")])
        for slide in prs2.slides
    )
    for forbidden in ["TODO", "state-of-the-art on all", "80% accuracy achieved", "mock result", "demo result"]:
        assert forbidden.lower() not in all_text.lower(), f"Forbidden phrase found: {forbidden}"
    measured_text = "\n".join(
        "\n".join([getattr(shape, "text", "") for shape in prs2.slides[idx].shapes if hasattr(shape, "text")])
        for idx in range(0, 32)
    )
    for target in ["75%", "80%", "85%"]:
        assert target not in measured_text, f"Target value {target} leaked before roadmap slide"
    assert "https://aclanthology.org/J05-1004/" in all_text
    assert "https://arxiv.org/abs/2305.14314" in all_text
    return {"slides": len(prs2.slides), "images_used": len(used_images)}


if __name__ == "__main__":
    build_deck()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    validation = validate_deck()
    print(f"Created: {OUT}")
    print(f"Slides: {validation['slides']}")
    print(f"Images used: {validation['images_used']}")
    print("Validation: passed")
