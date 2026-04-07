from __future__ import annotations

import json
from pathlib import Path
from typing import Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
PROJECT = ROOT / "srl_qa_project"
RESULTS = PROJECT / "results"
PLOTS = RESULTS / "plots"
OUTPUTS = PROJECT / "outputs"
ASSETS = ROOT / "ppt_assets"
CONTACT = ASSETS / "contact_sheets"
PPT_PATH = ROOT / "FINAL_PROJECT_PRESENTATION_40_SLIDES.pptx"

NAVY = RGBColor(13, 27, 42)
BLUE = RGBColor(31, 78, 121)
TEAL = RGBColor(27, 153, 139)
GOLD = RGBColor(217, 164, 65)
INK = RGBColor(34, 40, 49)
SLATE = RGBColor(94, 108, 132)
LIGHT = RGBColor(245, 248, 252)
PANEL = RGBColor(255, 255, 255)
BORDER = RGBColor(217, 223, 232)
SUCCESS = RGBColor(33, 128, 94)
WARN = RGBColor(191, 95, 36)


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def pct(value: float) -> str:
    return f"{value * 100:.1f}%"


def fmt(value: float, digits: int = 2) -> str:
    return f"{value:.{digits}f}"


class DeckBuilder:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.slide_index = 0
        self.title_font = "Aptos Display"
        self.body_font = "Aptos"

    def add_slide(self, title: str, section: str) -> object:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_index += 1
        self._background(slide)
        self._header(slide, title, section)
        self._footer(slide)
        return slide

    def _background(self, slide: object) -> None:
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = LIGHT

    def _header(self, slide: object, title: str, section: str) -> None:
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.6)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = NAVY
        header.line.fill.background()

        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.58), Inches(13.333), Inches(0.06)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = TEAL
        accent.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.45), Inches(0.12), Inches(8.9), Inches(0.33)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = self.title_font
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = PANEL

        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(10.9),
            Inches(0.12),
            Inches(1.95),
            Inches(0.3),
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = GOLD
        pill.line.fill.background()
        tf = pill.text_frame
        tf.clear()
        p2 = tf.paragraphs[0]
        p2.text = section.upper()
        p2.alignment = PP_ALIGN.CENTER
        p2.font.name = self.body_font
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = NAVY

    def _footer(self, slide: object) -> None:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(7.13), Inches(12.5), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER
        line.line.fill.background()

        footer = slide.shapes.add_textbox(
            Inches(0.45), Inches(7.15), Inches(11.5), Inches(0.18)
        )
        p = footer.text_frame.paragraphs[0]
        p.text = "PropQA-Net | SRL-Anchored QA | Final Research Presentation"
        p.font.name = self.body_font
        p.font.size = Pt(9)
        p.font.color.rgb = SLATE

        num = slide.shapes.add_textbox(
            Inches(12.1), Inches(7.12), Inches(0.7), Inches(0.2)
        )
        p2 = num.text_frame.paragraphs[0]
        p2.text = str(self.slide_index)
        p2.alignment = PP_ALIGN.RIGHT
        p2.font.name = self.body_font
        p2.font.size = Pt(10)
        p2.font.bold = True
        p2.font.color.rgb = SLATE

    def text_box(
        self,
        slide: object,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        size: int = 16,
        color: RGBColor = INK,
        bold: bool = False,
        align: PP_ALIGN = PP_ALIGN.LEFT,
    ) -> object:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.name = self.body_font
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        return box

    def bullet_box(
        self,
        slide: object,
        bullets: Sequence[str],
        left: float,
        top: float,
        width: float,
        height: float,
        size: int = 18,
        color: RGBColor = INK,
    ) -> object:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        for idx, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.bullet = True
            p.font.name = self.body_font
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.space_after = Pt(7)
        return box

    def card(
        self,
        slide: object,
        left: float,
        top: float,
        width: float,
        height: float,
        title: str,
        value: str,
        accent: RGBColor = TEAL,
    ) -> None:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL
        shape.line.color.rgb = BORDER

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(0.12),
            Inches(height),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        self.text_box(slide, title, left + 0.28, top + 0.18, width - 0.45, 0.28, 11, SLATE, True)
        self.text_box(slide, value, left + 0.28, top + 0.45, width - 0.45, 0.44, 21, NAVY, True)

    def ribbon(self, slide: object, text: str, left: float, top: float, width: float) -> None:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(0.28),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.body_font
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = PANEL

    def picture_contain(
        self, slide: object, path: Path, left: float, top: float, width: float, height: float
    ) -> None:
        if not path.exists():
            self.text_box(slide, f"Missing image:\n{path.name}", left, top, width, height, 14, WARN, True)
            return
        with Image.open(path) as img:
            img_w, img_h = img.size
        box_w = Inches(width)
        box_h = Inches(height)
        scale = min(box_w / img_w, box_h / img_h)
        pic_w = img_w * scale
        pic_h = img_h * scale
        pos_left = Inches(left) + (box_w - pic_w) / 2
        pos_top = Inches(top) + (box_h - pic_h) / 2

        frame = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        frame.fill.solid()
        frame.fill.fore_color.rgb = PANEL
        frame.line.color.rgb = BORDER
        slide.shapes.add_picture(str(path), pos_left, pos_top, width=pic_w, height=pic_h)

    def add_table(
        self,
        slide: object,
        data: Sequence[Sequence[str]],
        left: float,
        top: float,
        width: float,
        height: float,
        font_size: int = 12,
        header_fill: RGBColor = NAVY,
    ) -> None:
        rows = len(data)
        cols = len(data[0])
        table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                cell.text = str(data[r][c])
                cell.fill.solid()
                cell.fill.fore_color.rgb = PANEL if r else header_fill
                cell.text_frame.word_wrap = True
                for p in cell.text_frame.paragraphs:
                    p.font.name = self.body_font
                    p.font.size = Pt(font_size)
                    p.font.bold = r == 0
                    p.font.color.rgb = PANEL if r == 0 else INK
                    p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for c in range(cols):
            table.columns[c].width = Inches(width / cols)

    def add_notes(self, slide: object, text: str) -> None:
        notes = slide.notes_slide.notes_text_frame
        notes.text = text

    def save(self, path: Path) -> None:
        self.prs.save(str(path))


def build_deck() -> None:
    metrics = load_json(RESULTS / "metrics.json")
    stats = load_json(RESULTS / "data_statistics.json")
    bench = load_json(RESULTS / "benchmarks" / "benchmark_results.json")

    qa = metrics["qa_performance"]
    srl = metrics["srl_performance"]
    tracks = bench["tracks"]
    challenge = {
        name: tracks[name]["challenge"]
        for name in [
            "classical_baseline",
            "heuristic_reranker",
            "transformer_qa_assist",
            "full_hybrid",
        ]
    }
    test_subset = {
        name: tracks[name]["test_subset"]
        for name in [
            "classical_baseline",
            "heuristic_reranker",
            "transformer_qa_assist",
            "full_hybrid",
        ]
    }

    full_hybrid_challenge_samples = tracks["full_hybrid"]["records"]["challenge"]
    sample_loc = next(
        item for item in full_hybrid_challenge_samples if item["example_id"] == "challenge_005"
    )
    sample_mnr = next(
        item for item in full_hybrid_challenge_samples if item["example_id"] == "challenge_006"
    )
    sample_arg2 = next(
        item for item in full_hybrid_challenge_samples if item["example_id"] == "challenge_008"
    )

    ext_sources = [
        ("Automatic Labeling of Semantic Roles", "https://aclanthology.org/J02-3001/"),
        ("Large-Scale QA-SRL Parsing", "https://aclanthology.org/P18-1191/"),
        (
            "PropBank Comes of Age—Larger, Smarter, and more Diverse",
            "https://aclanthology.org/2022.starsem-1.24/",
        ),
        (
            "Potential and Limitations of LLMs in Capturing Structured Semantics",
            "https://arxiv.org/abs/2405.06410",
        ),
        (
            "LLMs Can Also Do Well! Breaking Barriers in Semantic Role Labeling via Large Language Models",
            "https://aclanthology.org/2025.findings-acl.1189/",
        ),
        (
            "BERT: Pre-training of Deep Bidirectional Transformers for Language Understanding",
            "https://arxiv.org/abs/1810.04805",
        ),
    ]

    d = DeckBuilder()

    slide = d.add_slide("Solution Snapshot", "Solution")
    d.text_box(
        slide,
        "Question Answering Using Semantic Roles",
        0.55,
        0.9,
        5.6,
        0.45,
        24,
        NAVY,
        True,
    )
    d.text_box(
        slide,
        "A reproducible SRL-QA system that combines PropBank-backed span supervision, a BiLSTM multi-task backbone, and role-aware hybrid inference.",
        0.55,
        1.35,
        5.7,
        0.9,
        16,
        INK,
    )
    d.bullet_box(
        slide,
        [
            "23,007 PropBank-derived QA pairs from 9,073 usable Treebank-backed instances",
            "Joint SRL + extractive QA baseline with 51.8% EM and 76.1% token F1",
            "Four-track benchmark suite comparing baseline, heuristic reranker, transformer assist, and full hybrid",
            "Research-grade outputs: plots, PDFs, implementation bundle, and Streamlit demo",
        ],
        0.6,
        2.35,
        5.7,
        3.3,
        18,
    )
    d.card(slide, 6.6, 0.95, 2.0, 1.0, "QA Exact Match", pct(qa["exact_match"]), TEAL)
    d.card(slide, 8.8, 0.95, 2.0, 1.0, "QA Token F1", pct(qa["token_f1"]), BLUE)
    d.card(slide, 11.0, 0.95, 1.8, 1.0, "SRL Micro F1", pct(srl["micro_f1"]), GOLD)
    d.picture_contain(slide, PLOTS / "research_architecture.png", 6.55, 2.1, 6.15, 4.6)
    d.add_notes(
        slide,
        "Open by positioning the work as a solution-first research system, not a generic NLP demo. "
        "Highlight the three layers: semantic supervision, multi-task baseline, and hybrid inference. "
        "Stress that the deck avoids a long problem intro and moves directly into what was built and measured. "
        "Animation suggestion: Fade in the three metric cards, then reveal the architecture image.",
    )

    slide = d.add_slide("End-to-End Research System", "Solution")
    d.bullet_box(
        slide,
        [
            "Local NLTK PropBank + Treebank assets ground the entire pipeline in real semantic annotations",
            "Data loader reconstructs spans, BIO labels, and natural-language questions deterministically",
            "PropQA-Net learns SRL tagging and answer extraction jointly",
            "Hybrid layer adds role-aware reranking, heuristics, and optional transformer candidates",
            "Benchmark runner, PDFs, and Streamlit app complete the research workflow",
        ],
        0.6,
        0.95,
        4.8,
        5.9,
        18,
    )
    d.picture_contain(slide, PLOTS / "research_architecture.png", 5.6, 1.05, 7.1, 5.7)
    d.add_notes(
        slide,
        "Walk the audience through the architecture left to right. "
        "Emphasize that the hybrid layer is an extension of the baseline rather than a replacement. "
        "This makes the benchmark comparisons meaningful and keeps the design interpretable. "
        "Transition suggestion: simple Fade.",
    )

    slide = d.add_slide("Repository Deliverables at a Glance", "Solution")
    d.card(slide, 0.7, 1.1, 2.7, 1.2, "Core Code", "main.py, data_loader.py, model.py, hybrid_qa.py", BLUE)
    d.card(slide, 3.6, 1.1, 2.8, 1.2, "Research Docs", "Survey, architecture, analysis, innovation", TEAL)
    d.card(slide, 6.7, 1.1, 2.7, 1.2, "Result Artifacts", "metrics.json, benchmark_results.json, plots/", GOLD)
    d.card(slide, 9.7, 1.1, 2.9, 1.2, "Submission Outputs", "survey.pdf, analysis.pdf, innovation.pdf, research_paper.pdf", SUCCESS)
    d.bullet_box(
        slide,
        [
            "The codebase is organized like a research package rather than a single training script",
            "Every major presentation claim is backed by a file, plot, JSON artifact, or PDF deliverable",
            "The same repo supports CLI execution, benchmarking, report generation, and interactive demo use",
        ],
        0.75,
        2.75,
        5.5,
        2.2,
        18,
    )
    d.picture_contain(slide, CONTACT / "research_paper_contact_sheet.png", 6.2, 2.2, 6.0, 4.3)
    d.add_notes(
        slide,
        "Use this slide to show maturity of the project package. "
        "The deck is backed by code, docs, metrics, plots, and generated reports. "
        "This supports the narrative that the project is reproducible and presentation-ready. "
        "Animation suggestion: stagger the four cards from left to right.",
    )

    slide = d.add_slide("Local Corpus Foundation and Filtering", "Solution")
    d.card(slide, 0.75, 1.0, 2.6, 1.1, "Visible PropBank Instances", f"{stats['total_propbank_instances']:,}", BLUE)
    d.card(slide, 3.55, 1.0, 2.6, 1.1, "Usable Instances", f"{stats['usable_propbank_instances']:,}", TEAL)
    d.card(slide, 6.35, 1.0, 2.4, 1.1, "QA Pairs", f"{stats['qa_pair_count']:,}", GOLD)
    d.card(slide, 8.95, 1.0, 1.75, 1.1, "Predicates", f"{stats['unique_predicates']:,}", SUCCESS)
    d.card(slide, 10.9, 1.0, 1.75, 1.1, "Rolesets", f"{stats['unique_rolesets']:,}", WARN)
    d.bullet_box(
        slide,
        [
            "Only Treebank-backed PropBank instances are retained so token spans can be reconstructed deterministically",
            "This filter trades raw corpus size for answer-span reliability and reproducibility",
            "The mean sentence length is 28.6 tokens; the mean answer length is 5.28 tokens",
            "The resulting corpus is small enough for a lightweight BiLSTM pipeline but large enough for structured analysis",
        ],
        0.8,
        2.45,
        5.8,
        3.6,
        17,
    )
    d.picture_contain(slide, PLOTS / "srl_pipeline.png", 7.0, 2.2, 5.5, 4.4)
    d.add_notes(
        slide,
        "Explain that the Treebank alignment filter is a deliberate quality step. "
        "The audience should see this as a methodological strength, not a limitation. "
        "Reliable span reconstruction matters because the project turns semantic arguments into supervised QA examples. "
        "Animation suggestion: bring in the cards first, then the pipeline figure.",
    )

    slide = d.add_slide("From PropBank Roles to QA Pairs", "Solution")
    d.bullet_box(
        slide,
        [
            "The loader reads real predicate-argument structures and converts them into (context, question, answer, SRL tags) tuples",
            "BIO labels preserve token-level semantic structure while answer spans support extractive QA training",
            "Template-based questions keep supervision deterministic and semantically aligned",
            "Example roles map naturally to WHO, WHAT, WHEN, WHERE, HOW, and WHY questions",
        ],
        0.6,
        1.0,
        5.2,
        3.5,
        17,
    )
    d.add_table(
        slide,
        [
            ["Role", "Question Pattern"],
            ["ARG0", "Who {predicate}?"],
            ["ARG1", "What did {subject} {predicate}?"],
            ["ARGM-TMP", "When did {subject} {predicate}?"],
            ["ARGM-LOC", "Where did {subject} {predicate}?"],
            ["ARGM-MNR", "How did {subject} {predicate}?"],
            ["ARGM-CAU", "Why did {subject} {predicate}?"],
        ],
        0.75,
        4.15,
        5.25,
        2.15,
        12,
    )
    d.picture_contain(slide, PLOTS / "propbank_example.png", 6.25, 1.0, 6.2, 5.9)
    d.add_notes(
        slide,
        "The key message is that the data pipeline is the conceptual bridge of the project. "
        "It turns semantic roles into question-answer supervision without inventing labels by hand. "
        "Mention that the deterministic question templates are intentional: they trade stylistic variety for clean supervision. "
        "Animation suggestion: fade in the example image after the question template table.",
    )

    slide = d.add_slide("Dataset Composition and Split Strategy", "Experiments")
    d.card(slide, 0.8, 1.0, 2.2, 1.0, "Train", f"{stats['split_sizes']['train']:,}", TEAL)
    d.card(slide, 3.25, 1.0, 2.2, 1.0, "Validation", f"{stats['split_sizes']['validation']:,}", BLUE)
    d.card(slide, 5.7, 1.0, 2.2, 1.0, "Test", f"{stats['split_sizes']['test']:,}", GOLD)
    d.card(slide, 8.15, 1.0, 2.2, 1.0, "Question Types", "6", SUCCESS)
    d.card(slide, 10.6, 1.0, 2.2, 1.0, "Seed", "42", WARN)
    d.bullet_box(
        slide,
        [
            f"WHAT dominates the corpus with {stats['qa_pairs_per_question_type']['WHAT']:,} pairs; WHO follows with {stats['qa_pairs_per_question_type']['WHO']:,}",
            f"WHEN ({stats['qa_pairs_per_question_type']['WHEN']:,}) and HOW ({stats['qa_pairs_per_question_type']['HOW']:,}) are mid-frequency categories",
            f"WHERE ({stats['qa_pairs_per_question_type']['WHERE']:,}) and WHY ({stats['qa_pairs_per_question_type']['WHY']:,}) are relatively sparse",
            "This imbalance directly influences exact-match behavior and long-tail SRL performance",
        ],
        0.8,
        2.3,
        5.55,
        3.7,
        17,
    )
    d.picture_contain(slide, PLOTS / "dataset_balance.png", 6.45, 1.9, 6.0, 4.9)
    d.add_notes(
        slide,
        "Use this slide to connect label distribution with later performance. "
        "The scarcity of WHY and some modifier roles is one reason the hybrid extension matters. "
        "This also explains why the project reports both micro and macro metrics later on. "
        "Animation suggestion: one group fade for cards, one for the plot.",
    )

    slide = d.add_slide("Span and Sequence Complexity Profile", "Experiments")
    d.bullet_box(
        slide,
        [
            "Average answer spans are short enough for extractive decoding, but the long tail still includes answers up to 62 tokens",
            "Sentence length ranges from 3 to 249 tokens, which justifies the max-length caps in config.py",
            "The model caps context at 128 tokens and questions at 32 tokens for stable batching",
            "Many boundary errors come from semantically valid but overly long or overly short span choices",
        ],
        0.65,
        1.0,
        5.5,
        3.0,
        17,
    )
    d.card(slide, 0.8, 4.4, 2.5, 1.0, "Mean Sentence Length", fmt(stats["sentence_length_summary"]["mean"], 1), BLUE)
    d.card(slide, 3.55, 4.4, 2.5, 1.0, "Mean Answer Length", fmt(stats["answer_length_summary"]["mean"], 2), TEAL)
    d.card(slide, 0.8, 5.55, 2.5, 1.0, "Sentence Max", str(int(stats["sentence_length_summary"]["max"])), GOLD)
    d.card(slide, 3.55, 5.55, 2.5, 1.0, "Answer Max", str(int(stats["answer_length_summary"]["max"])), WARN)
    d.picture_contain(slide, PLOTS / "answer_length_dist.png", 6.4, 1.0, 6.1, 5.9)
    d.add_notes(
        slide,
        "Explain that the dataset is extractive-friendly overall, but not trivial. "
        "The long-tail answer distribution is important because it explains why token F1 can be good even when exact match falls. "
        "This also motivates the hybrid reranker's explicit role and shape features. "
        "Transition suggestion: gentle Fade only.",
    )

    slide = d.add_slide("PropQA-Net Baseline Architecture", "Solution")
    d.bullet_box(
        slide,
        [
            "Context encoder: shared word embeddings + POS embeddings + predicate flags into a BiLSTM",
            "Question encoder: shared word embeddings into a second BiLSTM followed by masked mean pooling",
            "Two heads: token-level SRL classifier and answer boundary predictors",
            "Decoding combines BIO-derived candidate spans with question alignment and boundary confidence",
        ],
        0.65,
        1.0,
        5.2,
        3.2,
        17,
    )
    d.add_table(
        slide,
        [
            ["Hyperparameter", "Value"],
            ["Word embedding dim", "100"],
            ["POS embedding dim", "32"],
            ["Predicate embedding dim", "8"],
            ["Hidden size", "128"],
            ["Dropout", "0.30"],
            ["Alpha", "0.50"],
        ],
        0.8,
        4.2,
        4.2,
        2.15,
        12,
    )
    d.picture_contain(slide, PLOTS / "propqa_architecture.png", 5.45, 1.0, 7.0, 5.9)
    d.add_notes(
        slide,
        "Describe this as a deliberate middle path between classical features and heavy transformers. "
        "The predicate flag is especially important because it anchors the semantic event of interest. "
        "Point out that the architecture is lightweight enough for offline reproducibility and rich enough for multi-task learning. "
        "Animation suggestion: fade in the architecture after the left-side explanation.",
    )

    slide = d.add_slide("Learning Objective and Decode Logic", "Solution")
    d.text_box(slide, "Multi-task training objective", 0.75, 1.0, 4.0, 0.25, 18, NAVY, True)
    d.text_box(slide, "L = alpha * L_SRL + (1 - alpha) * L_QA", 0.85, 1.45, 4.2, 0.3, 22, BLUE, True)
    d.text_box(slide, "L_SRL = token-level cross-entropy over BIO labels\nL_QA = average of start and end cross-entropy", 0.85, 1.95, 4.6, 0.8, 17, INK)
    d.text_box(slide, "Decode-time fusion", 0.75, 3.0, 4.0, 0.25, 18, NAVY, True)
    d.text_box(slide, "score = 0.60 * cosine(span, question) + 0.40 * boundary_confidence", 0.85, 3.45, 4.3, 0.45, 18, TEAL, True)
    d.bullet_box(
        slide,
        [
            "BIO decoding yields semantically plausible candidate spans",
            "Boundary scoring keeps answer extraction compatible with standard QA practice",
            "The fallback mechanism protects against malformed BIO sequences",
        ],
        0.8,
        4.15,
        4.8,
        1.8,
        16,
    )
    d.picture_contain(slide, PLOTS / "hybridpropqa.png", 5.7, 1.0, 6.6, 5.8)
    d.add_notes(
        slide,
        "This slide is the conceptual center of the model. "
        "Explain that SRL and QA are coupled both during training and during decoding. "
        "The model does not only predict spans; it predicts spans in the context of semantic-role structure. "
        "Animation suggestion: reveal the formulas first, then the candidate-selection bullets.",
    )

    slide = d.add_slide("Training Setup and Reproducibility", "Experiments")
    d.card(slide, 0.8, 1.0, 2.0, 1.0, "Batch Size", "64", BLUE)
    d.card(slide, 3.0, 1.0, 2.0, 1.0, "Max Epochs", "6", TEAL)
    d.card(slide, 5.2, 1.0, 2.0, 1.0, "Patience", "5", GOLD)
    d.card(slide, 7.4, 1.0, 2.0, 1.0, "Learning Rate", "1e-3", SUCCESS)
    d.card(slide, 9.6, 1.0, 2.0, 1.0, "Weight Decay", "1e-5", WARN)
    d.bullet_box(
        slide,
        [
            "config.py centralizes all paths and hyperparameters",
            "main.py supports train, eval, infer, ask, benchmark, report, and app modes",
            "The pipeline saves checkpoints, metrics, plots, and report artifacts automatically",
            "The repo includes local corpus assets and generated PDFs, which supports offline reruns",
        ],
        0.8,
        2.35,
        5.25,
        3.5,
        17,
    )
    d.picture_contain(slide, PLOTS / "loss_curve.png", 6.15, 1.9, 6.2, 4.9)
    d.add_notes(
        slide,
        "Frame the training setup as intentionally conservative and reproducible. "
        "The point is not raw benchmark chasing but a stable, inspectable project pipeline. "
        "Mention the saved checkpoint and the report mode as evidence that the repo is more than a notebook experiment. "
        "Animation suggestion: wipe in the loss curve from right.",
    )

    slide = d.add_slide("Experiment Matrix", "Experiments")
    d.add_table(
        slide,
        [
            ["Evaluation Component", "Setting"],
            ["Full held-out test set", f"{stats['split_sizes']['test']:,} examples"],
            ["Challenge suite", str(bench["metadata"]["challenge_size"]) + " curated role-sensitive examples"],
            ["Benchmark test subset", str(bench["metadata"]["test_subset_size"]) + " sampled examples"],
            ["Primary QA metrics", "Exact Match, Token F1"],
            ["Primary SRL metrics", "Micro F1, Macro F1, BIO Accuracy"],
            ["Tracks compared", "Baseline, Heuristic, Transformer Assist, Full Hybrid"],
        ],
        0.75,
        1.0,
        5.7,
        4.2,
        12,
    )
    d.bullet_box(
        slide,
        [
            "The benchmark separates authoritative repository metrics from hybrid ablation analysis",
            "Challenge evaluation stresses role-sensitive questions like location, manner, cause, and recipient extraction",
            "The test-subset benchmark measures whether those gains generalize beyond the curated challenge set",
        ],
        0.9,
        5.35,
        5.45,
        1.6,
        15,
    )
    d.picture_contain(slide, PLOTS / "ablation_summary.png", 6.55, 1.15, 5.85, 5.6)
    d.add_notes(
        slide,
        "Clarify that there are two result regimes: the main repository metrics and the hybrid ablation suite. "
        "This separation prevents confusion between baseline training results and engineering experiments on top. "
        "Point out that the challenge set is targeted, not a replacement for the full test data. "
        "Animation suggestion: table first, then chart.",
    )

    slide = d.add_slide("Baseline QA Results", "Results")
    d.card(slide, 0.8, 1.0, 2.2, 1.0, "QA Exact Match", pct(qa["exact_match"]), BLUE)
    d.card(slide, 3.25, 1.0, 2.2, 1.0, "QA Token F1", pct(qa["token_f1"]), TEAL)
    d.card(slide, 5.7, 1.0, 2.2, 1.0, "Length Deviation", fmt(qa["answer_length_deviation_mean"], 2), GOLD)
    d.card(slide, 8.15, 1.0, 2.2, 1.0, "Abs. Length Dev.", fmt(qa["answer_length_deviation_abs_mean"], 2), SUCCESS)
    d.bullet_box(
        slide,
        [
            "WHO and WHAT are the strongest question families by exact match and token F1",
            "WHEN remains relatively manageable because temporal expressions are often compact and distinctive",
            "WHY has the weakest exact match, indicating frequent partial-span or boundary mismatches",
            "The baseline is strongest on common semantic roles and weaker on the sparse long tail",
        ],
        0.8,
        2.35,
        5.1,
        3.8,
        17,
    )
    d.picture_contain(slide, PLOTS / "qa_accuracy_by_qtype.png", 6.2, 1.8, 6.2, 4.9)
    d.add_notes(
        slide,
        "Tell the audience that these are the authoritative local baseline metrics from metrics.json. "
        "Use the question-type plot to explain why the project later invests in role-aware hybrid improvements. "
        "Boundary sensitivity is especially visible for WHY and some longer WHAT answers. "
        "Animation suggestion: subtle float-in for the metric cards.",
    )

    slide = d.add_slide("Baseline SRL Results", "Results")
    d.card(slide, 0.8, 1.0, 2.3, 1.0, "SRL Micro F1", pct(srl["micro_f1"]), BLUE)
    d.card(slide, 3.35, 1.0, 2.3, 1.0, "SRL Macro F1", pct(srl["macro_f1"]), TEAL)
    d.card(slide, 5.9, 1.0, 2.3, 1.0, "BIO Accuracy", pct(srl["bio_accuracy"]), GOLD)
    d.bullet_box(
        slide,
        [
            "ARG1 and ARG0 dominate performance because they are frequent and structurally stable",
            "Modifier roles like ARGM-MOD and ARGM-NEG score very highly when they appear",
            "Macro F1 remains low because many specialized labels have extremely small support",
            "This gap between micro and macro is a clear long-tail effect rather than a contradiction in the model",
        ],
        0.8,
        2.35,
        5.1,
        3.9,
        17,
    )
    d.picture_contain(slide, PLOTS / "f1_by_argtype.png", 6.1, 1.7, 6.3, 5.0)
    d.add_notes(
        slide,
        "Explain micro versus macro carefully. "
        "Micro captures overall usefulness on the frequent roles; macro exposes the difficulty of the rare label tail. "
        "This slide supports an honest and credible discussion of baseline quality. "
        "Animation suggestion: chart fade only.",
    )

    slide = d.add_slide("Per-Question and Per-Role Analysis", "Analysis")
    d.add_table(
        slide,
        [
            ["Question Type", "EM", "F1", "Count"],
            ["WHO", pct(qa["per_question_type"]["WHO"]["em"]), pct(qa["per_question_type"]["WHO"]["f1"]), str(int(qa["per_question_type"]["WHO"]["count"]))],
            ["WHAT", pct(qa["per_question_type"]["WHAT"]["em"]), pct(qa["per_question_type"]["WHAT"]["f1"]), str(int(qa["per_question_type"]["WHAT"]["count"]))],
            ["WHEN", pct(qa["per_question_type"]["WHEN"]["em"]), pct(qa["per_question_type"]["WHEN"]["f1"]), str(int(qa["per_question_type"]["WHEN"]["count"]))],
            ["WHERE", pct(qa["per_question_type"]["WHERE"]["em"]), pct(qa["per_question_type"]["WHERE"]["f1"]), str(int(qa["per_question_type"]["WHERE"]["count"]))],
            ["HOW", pct(qa["per_question_type"]["HOW"]["em"]), pct(qa["per_question_type"]["HOW"]["f1"]), str(int(qa["per_question_type"]["HOW"]["count"]))],
            ["WHY", pct(qa["per_question_type"]["WHY"]["em"]), pct(qa["per_question_type"]["WHY"]["f1"]), str(int(qa["per_question_type"]["WHY"]["count"]))],
        ],
        0.75,
        1.0,
        5.25,
        3.25,
        11,
    )
    d.add_table(
        slide,
        [
            ["Role", "F1", "Support"],
            ["ARG1", pct(metrics["srl_performance"]["per_role"]["ARG1"]["f1"]), str(metrics["srl_performance"]["per_role"]["ARG1"]["support"])],
            ["ARG0", pct(metrics["srl_performance"]["per_role"]["ARG0"]["f1"]), str(metrics["srl_performance"]["per_role"]["ARG0"]["support"])],
            ["ARGM-CAU", pct(metrics["srl_performance"]["per_role"]["ARGM-CAU"]["f1"]), str(metrics["srl_performance"]["per_role"]["ARGM-CAU"]["support"])],
            ["ARGM-TMP", pct(metrics["srl_performance"]["per_role"]["ARGM-TMP"]["f1"]), str(metrics["srl_performance"]["per_role"]["ARGM-TMP"]["support"])],
            ["ARGM-LOC", pct(metrics["srl_performance"]["per_role"]["ARGM-LOC"]["f1"]), str(metrics["srl_performance"]["per_role"]["ARGM-LOC"]["support"])],
            ["ARGM-MNR", pct(metrics["srl_performance"]["per_role"]["ARGM-MNR"]["f1"]), str(metrics["srl_performance"]["per_role"]["ARGM-MNR"]["support"])],
        ],
        0.75,
        4.55,
        5.25,
        2.0,
        11,
    )
    d.bullet_box(
        slide,
        [
            "The baseline is semantically useful, but not uniformly strong across every role or question type",
            "WHY and some modifier-heavy categories are the clearest opportunities for targeted improvements",
            "This analysis directly motivates the hybrid role-aware benchmark extension",
        ],
        6.2,
        1.1,
        6.0,
        2.0,
        18,
    )
    d.picture_contain(slide, PLOTS / "question_type_heatmap.png", 6.15, 3.0, 6.15, 3.6)
    d.add_notes(
        slide,
        "Use the tables for precise reporting and the heatmap for intuition. "
        "This slide should set up the story that the hybrid system is targeting specific weaknesses rather than replacing a bad baseline. "
        "Animation suggestion: bring in the two tables together, then the heatmap.",
    )

    slide = d.add_slide("Confusion Structure", "Analysis")
    d.bullet_box(
        slide,
        [
            "Frequent roles dominate the prediction landscape, which is expected in PropBank-style corpora",
            "Confusions often occur when a plausible span fits the sentence but not the intended semantic role",
            "Object-like, location-like, and recipient-like spans are especially important for later reranking",
            "This plot motivates explicit role-match features rather than relying on boundary confidence alone",
        ],
        0.75,
        1.0,
        4.8,
        3.6,
        17,
    )
    d.picture_contain(slide, PLOTS / "confusion_matrix.png", 5.55, 0.95, 7.0, 5.9)
    d.add_notes(
        slide,
        "Explain that the confusion matrix is not only a diagnostic; it also shapes the hybrid design. "
        "A high-confidence span can still be semantically wrong, so role-aware reranking is a natural extension. "
        "Animation suggestion: no complex transitions; use a static display for readability.",
    )

    slide = d.add_slide("Error Taxonomy", "Analysis")
    d.bullet_box(
        slide,
        [
            "The project separates errors into predicate misses, wrong-role predictions, span-boundary mistakes, and other residual cases",
            "Boundary mistakes are particularly important because they hurt exact match more than token F1",
            "Wrong-role errors matter most for semantically sensitive questions such as WHERE, HOW, and TO-WHOM",
            "The hybrid extension is designed to suppress exactly these role-sensitive mistakes",
        ],
        0.75,
        1.0,
        4.9,
        3.5,
        17,
    )
    d.picture_contain(slide, PLOTS / "error_taxonomy.png", 5.55, 1.15, 6.9, 5.4)
    d.add_notes(
        slide,
        "Make the distinction between semantic plausibility and exact correctness. "
        "A model can choose a good region of text and still be penalized on exact match because of span boundaries. "
        "That nuance is why the deck keeps both EM and token F1 in view. "
        "Animation suggestion: Fade only.",
    )

    slide = d.add_slide("Qualitative Failure Cases", "Analysis")
    d.bullet_box(
        slide,
        [
            "Long WHAT answers can be truncated or shifted to adjacent spans",
            "Manner and location questions can be answered with larger ARG1-like spans when role cues are weak",
            "Temporal and causal phrases often receive partial credit under token F1 but fail exact match",
            "These failure patterns are visible in both the saved prediction sample and the error gallery figure",
        ],
        0.75,
        1.0,
        4.9,
        3.6,
        17,
    )
    d.picture_contain(slide, PLOTS / "error_gallery.png", 5.55, 1.0, 6.9, 5.8)
    d.add_notes(
        slide,
        "Use one or two examples verbally while the gallery is visible. "
        "This slide makes the analysis concrete and prepares the audience to appreciate the benchmark examples later. "
        "Animation suggestion: none; keep the slide readable.",
    )

    slide = d.add_slide("Challenge Benchmark Results", "Results")
    d.add_table(
        slide,
        [
            ["Track", "EM", "F1", "Role Acc.", "Latency (ms)"],
            ["Classical baseline", pct(challenge["classical_baseline"]["exact_match"]), pct(challenge["classical_baseline"]["token_f1"]), pct(challenge["classical_baseline"]["role_accuracy"]), fmt(challenge["classical_baseline"]["mean_latency_ms"], 2)],
            ["Heuristic reranker", pct(challenge["heuristic_reranker"]["exact_match"]), pct(challenge["heuristic_reranker"]["token_f1"]), pct(challenge["heuristic_reranker"]["role_accuracy"]), fmt(challenge["heuristic_reranker"]["mean_latency_ms"], 2)],
            ["Transformer assist", pct(challenge["transformer_qa_assist"]["exact_match"]), pct(challenge["transformer_qa_assist"]["token_f1"]), pct(challenge["transformer_qa_assist"]["role_accuracy"]), fmt(challenge["transformer_qa_assist"]["mean_latency_ms"], 2)],
            ["Full hybrid", pct(challenge["full_hybrid"]["exact_match"]), pct(challenge["full_hybrid"]["token_f1"]), pct(challenge["full_hybrid"]["role_accuracy"]), fmt(challenge["full_hybrid"]["mean_latency_ms"], 2)],
        ],
        0.75,
        1.05,
        5.1,
        2.35,
        11,
    )
    d.bullet_box(
        slide,
        [
            "Role-aware upgrades lift challenge-set role accuracy from 20.0% to 100.0%",
            "The heuristic reranker already captures most of the measurable gain on this curated suite",
            "Transformer support does not improve challenge accuracy further in this environment, but it adds large latency costs",
        ],
        0.8,
        3.75,
        5.0,
        2.2,
        16,
    )
    d.picture_contain(slide, PLOTS / "challenge_table.png", 6.1, 1.0, 6.2, 5.9)
    d.add_notes(
        slide,
        "This is one of the strongest result slides. "
        "Frame it as evidence that role-aware reasoning changes answer behavior substantially on targeted semantic questions. "
        "Also point out that the cheapest upgraded path already captures the main gain. "
        "Animation suggestion: table first, then figure.",
    )

    slide = d.add_slide("Sampled Test-Subset Benchmark", "Results")
    d.add_table(
        slide,
        [
            ["Track", "EM", "F1", "Role Acc.", "Latency (ms)"],
            ["Classical baseline", pct(test_subset["classical_baseline"]["exact_match"]), pct(test_subset["classical_baseline"]["token_f1"]), pct(test_subset["classical_baseline"]["role_accuracy"]), fmt(test_subset["classical_baseline"]["mean_latency_ms"], 2)],
            ["Heuristic reranker", pct(test_subset["heuristic_reranker"]["exact_match"]), pct(test_subset["heuristic_reranker"]["token_f1"]), pct(test_subset["heuristic_reranker"]["role_accuracy"]), fmt(test_subset["heuristic_reranker"]["mean_latency_ms"], 2)],
            ["Transformer assist", pct(test_subset["transformer_qa_assist"]["exact_match"]), pct(test_subset["transformer_qa_assist"]["token_f1"]), pct(test_subset["transformer_qa_assist"]["role_accuracy"]), fmt(test_subset["transformer_qa_assist"]["mean_latency_ms"], 2)],
            ["Full hybrid", pct(test_subset["full_hybrid"]["exact_match"]), pct(test_subset["full_hybrid"]["token_f1"]), pct(test_subset["full_hybrid"]["role_accuracy"]), fmt(test_subset["full_hybrid"]["mean_latency_ms"], 2)],
        ],
        0.75,
        1.05,
        5.15,
        2.35,
        11,
    )
    d.bullet_box(
        slide,
        [
            "The hybrid tracks still outperform the baseline on the sampled test subset, but the gains are smaller than on the curated challenge set",
            "This suggests that the extension is especially effective on role-sensitive cases rather than universally dominant",
            "The same-result pattern across the three upgraded tracks reinforces the value of the lightweight heuristic reranker",
        ],
        0.8,
        3.8,
        5.0,
        2.15,
        16,
    )
    d.picture_contain(slide, PLOTS / "benchmark.png", 6.1, 1.05, 6.2, 5.8)
    d.add_notes(
        slide,
        "Present this slide honestly: the hybrid extension helps, but not uniformly or magically. "
        "This strengthens the credibility of the project because it reports trade-offs rather than only best cases. "
        "Animation suggestion: simple fade.",
    )

    slide = d.add_slide("Latency–Accuracy Trade-off", "Results")
    d.bullet_box(
        slide,
        [
            "Heuristic reranking is only slightly slower than the classical baseline",
            "Transformer-assisted and full-hybrid modes are much slower on CPU-bound local inference",
            "On the observed benchmark subsets, the heavier tracks do not outperform the heuristic reranker",
            "For classroom demos and practical offline use, the heuristic reranker is the best cost–benefit choice",
        ],
        0.75,
        1.05,
        5.15,
        3.5,
        17,
    )
    d.card(slide, 0.9, 4.9, 2.2, 1.0, "Best Practical Upgrade", "Heuristic reranker", TEAL)
    d.card(slide, 3.35, 4.9, 2.2, 1.0, "Strongest Challenge Role Acc.", "100.0%", BLUE)
    d.picture_contain(slide, PLOTS / "latency_accuracy_tradeoff.png", 6.0, 1.05, 6.35, 5.8)
    d.add_notes(
        slide,
        "Frame this slide as an engineering result. "
        "The project is not only about accuracy; it is also about what the system costs to run. "
        "This is particularly relevant for a final-year project that must work reliably on local hardware. "
        "Animation suggestion: cards then chart.",
    )

    slide = d.add_slide("Challenge Performance by Question Type", "Analysis")
    d.bullet_box(
        slide,
        [
            "Role-aware reranking is especially valuable for HOW, WHERE, and TO-WHOM questions",
            "These categories are exactly where role identity matters more than simple lexical overlap",
            "The heatmap shows that targeted semantic cues provide more benefit than generic span scoring alone",
        ],
        0.75,
        1.1,
        4.9,
        2.4,
        18,
    )
    d.picture_contain(slide, PLOTS / "question_type_heatmap.png", 5.7, 1.0, 6.5, 5.8)
    d.add_notes(
        slide,
        "Use this slide to explain why a semantic-role-aware system can outperform a pure span reader on certain categories. "
        "The benefit is not random; it is concentrated where role expectation matters. "
        "Animation suggestion: no motion beyond default fade.",
    )

    slide = d.add_slide("Challenge Performance by Target Role", "Analysis")
    d.bullet_box(
        slide,
        [
            "Role accuracy rises most strongly for roles that align with explicit heuristic patterns and question intent",
            "ARGM-LOC, ARGM-MNR, and ARG2 are especially important proof points for the hybrid design",
            "The heatmap reinforces that semantic expectations drive better candidate selection",
        ],
        0.75,
        1.1,
        4.9,
        2.4,
        18,
    )
    d.picture_contain(slide, PLOTS / "role_heatmap.png", 5.7, 1.0, 6.5, 5.8)
    d.add_notes(
        slide,
        "Tie this slide directly back to the reranker's role-match feature. "
        "The audience should see that the hybrid system is structurally motivated, not just a bundle of extra tricks. "
        "Animation suggestion: static display.",
    )

    slide = d.add_slide("Hybrid Reasoning Example: Location", "Analysis")
    d.ribbon(slide, "Challenge Example: WHERE / ARGM-LOC", 0.8, 0.95, 3.3)
    d.text_box(slide, f"Context: {sample_loc['context']}", 0.8, 1.45, 5.4, 1.2, 15, INK, False)
    d.text_box(slide, f"Question: {sample_loc['question']}", 0.8, 2.6, 5.4, 0.45, 17, NAVY, True)
    d.card(slide, 0.85, 3.25, 2.2, 1.0, "Expected Answer", sample_loc["expected_answer"], TEAL)
    d.card(slide, 3.25, 3.25, 2.2, 1.0, "Predicted Role", sample_loc["predicted_role"], BLUE)
    d.card(slide, 0.85, 4.45, 2.2, 1.0, "Exact Match", pct(sample_loc["exact_match"]), GOLD)
    d.card(slide, 3.25, 4.45, 2.2, 1.0, "Token F1", pct(sample_loc["token_f1"]), SUCCESS)
    d.text_box(slide, "Reasoning summary", 0.85, 5.7, 2.0, 0.25, 15, NAVY, True)
    d.text_box(slide, sample_loc["reasoning_summary"], 0.85, 6.0, 5.2, 0.7, 12, INK)
    d.add_table(
        slide,
        [["Candidate", "Role", "Source", "Score"]]
        + [[ev["text"], ev["role"], ev["source"], fmt(ev["score"], 3)] for ev in sample_loc["evidence_spans"]],
        6.35,
        1.25,
        6.0,
        4.9,
        11,
    )
    d.add_notes(
        slide,
        "This is a clean demonstration of the hybrid value proposition. "
        "The location heuristic proposes a role-compatible span and the reranker selects it over the baseline's larger ARG1-like answer. "
        "Use this as a live-demo backup example if needed. "
        "Animation suggestion: show the candidate table after summarizing the question.",
    )

    slide = d.add_slide("Hybrid Reasoning Example: Manner and Recipient", "Analysis")
    d.ribbon(slide, "Challenge Example: HOW / ARGM-MNR", 0.8, 0.95, 3.3)
    d.text_box(slide, f"Context: {sample_mnr['context']}", 0.8, 1.35, 5.0, 0.85, 14, INK)
    d.text_box(slide, f"Question: {sample_mnr['question']}", 0.8, 2.15, 5.0, 0.4, 16, NAVY, True)
    d.add_table(
        slide,
        [["Candidate", "Role", "Source", "Score"]]
        + [[ev["text"], ev["role"], ev["source"], fmt(ev["score"], 3)] for ev in sample_mnr["evidence_spans"]],
        0.8,
        2.75,
        5.25,
        2.6,
        10,
    )
    d.text_box(slide, "Recipient example", 6.35, 1.0, 2.0, 0.3, 16, NAVY, True)
    d.text_box(slide, f"Context: {sample_arg2['context']}", 6.35, 1.35, 5.5, 0.75, 14, INK)
    d.text_box(slide, f"Question: {sample_arg2['question']}", 6.35, 2.1, 5.4, 0.4, 16, NAVY, True)
    d.add_table(
        slide,
        [["Candidate", "Role", "Source", "Score"]]
        + [[ev["text"], ev["role"], ev["source"], fmt(ev["score"], 3)] for ev in sample_arg2["evidence_spans"]],
        6.35,
        2.65,
        5.25,
        2.7,
        10,
    )
    d.text_box(slide, "Both examples show the same pattern: question intent narrows the expected role, heuristics create role-aligned candidates, and weighted reranking chooses the best semantic match.", 0.8, 5.7, 11.0, 0.7, 14, INK, False)
    d.add_notes(
        slide,
        "Use this slide to illustrate that the hybrid logic generalizes beyond one case. "
        "The manner example shows how adverbial cues help; the recipient example shows how transfer-verb structure helps ARG2 selection. "
        "Animation suggestion: reveal the right panel after covering the left example.",
    )

    slide = d.add_slide("Why the Heuristic Reranker Matters", "Innovation")
    d.picture_contain(slide, PLOTS / "ablation_summary.png", 6.15, 1.0, 6.1, 4.6)
    d.bullet_box(
        slide,
        [
            "The heuristic reranker captures nearly all measurable gains observed on the challenge and sampled test benchmarks",
            "It stays fast enough for interactive local use, unlike heavier transformer-assisted modes",
            "The result is a strong argument for semantic feature design rather than only adding bigger models",
        ],
        0.8,
        1.05,
        5.0,
        2.4,
        18,
    )
    d.card(slide, 0.9, 4.0, 2.15, 1.0, "Challenge F1", pct(challenge["heuristic_reranker"]["token_f1"]), TEAL)
    d.card(slide, 3.25, 4.0, 2.15, 1.0, "Challenge Role Acc.", pct(challenge["heuristic_reranker"]["role_accuracy"]), BLUE)
    d.card(slide, 0.9, 5.2, 2.15, 1.0, "Subset F1", pct(test_subset["heuristic_reranker"]["token_f1"]), GOLD)
    d.card(slide, 3.25, 5.2, 2.15, 1.0, "Latency", fmt(test_subset["heuristic_reranker"]["mean_latency_ms"], 2) + " ms", SUCCESS)
    d.add_notes(
        slide,
        "This slide makes the engineering case for the project. "
        "The strongest practical improvement comes from a transparent heuristic-semantic layer, not from the heaviest model. "
        "That is a meaningful contribution because it changes what a student-scale project can do locally. "
        "Animation suggestion: cards rise in after the bullets.",
    )

    slide = d.add_slide("Innovation Summary", "Innovation")
    d.card(slide, 0.85, 1.1, 2.75, 1.3, "Innovation 1", "SRL-anchored extractive QA with joint learning", TEAL)
    d.card(slide, 3.95, 1.1, 2.75, 1.3, "Innovation 2", "Role-aware question parsing and intent mapping", BLUE)
    d.card(slide, 7.05, 1.1, 2.75, 1.3, "Innovation 3", "Multi-channel candidates: baseline + heuristics + transformer", GOLD)
    d.card(slide, 10.15, 1.1, 2.3, 1.3, "Innovation 4", "Transparent weighted semantic reranking", SUCCESS)
    d.bullet_box(
        slide,
        [
            "The project innovates most strongly through integration, transparency, and research packaging rather than through a single giant model claim",
            "Each innovation is visible in code, measurable in artifacts, and presentable in the app and PDFs",
            "This makes the system defensible in both technical and academic evaluation settings",
        ],
        0.9,
        3.0,
        5.2,
        2.2,
        17,
    )
    d.picture_contain(slide, PLOTS / "research_architecture.png", 6.2, 2.55, 6.05, 4.0)
    d.add_notes(
        slide,
        "Position innovation as system-level coherence. "
        "The deck should not imply that every subcomponent is individually unprecedented; the strength is in how they work together. "
        "Animation suggestion: reveal the four cards in sequence.",
    )

    slide = d.add_slide("Innovation Deep Dive: Role-Aware Parsing", "Innovation")
    d.add_table(
        slide,
        [
            ["Question Pattern", "Mapped Type", "Expected Role"],
            ["Who ...", "WHO", "ARG0"],
            ["What ...", "WHAT", "ARG1"],
            ["When ...", "WHEN", "ARGM-TMP"],
            ["Where ...", "WHERE", "ARGM-LOC"],
            ["How ...", "HOW", "ARGM-MNR"],
            ["Why ...", "WHY", "ARGM-CAU"],
            ["Whom / Who received ...", "TO-WHOM", "ARG2"],
        ],
        0.75,
        1.05,
        5.25,
        4.7,
        11,
    )
    d.bullet_box(
        slide,
        [
            "Question intent is parsed before final answer selection",
            "Expected roles become a reusable signal for filtering and ranking candidates",
            "This bridges the surface form of the question with the semantic structure of the answer",
        ],
        0.85,
        5.95,
        5.15,
        0.7,
        15,
    )
    d.picture_contain(slide, PLOTS / "hybridpropqa.png", 6.2, 1.05, 6.05, 5.7)
    d.add_notes(
        slide,
        "Explain that role-aware parsing is a lightweight but powerful design choice. "
        "It is one of the cleanest bridges between structured semantics and question answering behavior. "
        "Animation suggestion: no special effects beyond a fade.",
    )

    slide = d.add_slide("Innovation Deep Dive: Transparent Reranking", "Innovation")
    d.add_table(
        slide,
        [
            ["Feature", "Weight", "Purpose"],
            ["Base score", "0.30", "Respect source confidence"],
            ["Role match", "0.32", "Prioritize semantic fit to the question"],
            ["Semantic alignment", "0.22", "Use sentence-level compatibility"],
            ["Lexical overlap", "0.06", "Reward useful content overlap"],
            ["Shape bonus", "0.10", "Prefer role-appropriate span forms"],
            ["Baseline bonus", "dynamic", "Preserve useful model agreement"],
        ],
        0.75,
        1.05,
        5.55,
        4.4,
        11,
    )
    d.bullet_box(
        slide,
        [
            "The weights are visible in code, easy to explain, and easy to tune",
            "The reranker is interpretable enough to support deterministic reasoning traces",
            "This is more presentation-friendly and research-friendly than a hidden black-box reranker",
        ],
        0.85,
        5.7,
        5.3,
        0.9,
        15,
    )
    d.picture_contain(slide, PLOTS / "latency_accuracy_tradeoff.png", 6.55, 1.15, 5.75, 5.45)
    d.add_notes(
        slide,
        "Stress that transparent scoring is a design choice, not a compromise. "
        "It lets the project explain why an answer won, not only that it won. "
        "Animation suggestion: one-stage reveal for the table, then the chart.",
    )

    slide = d.add_slide("Innovation Deep Dive: Research Packaging", "Innovation")
    d.card(slide, 0.8, 1.0, 2.4, 1.1, "Benchmarking", "Four-track comparison with challenge + test subsets", BLUE)
    d.card(slide, 3.45, 1.0, 2.4, 1.1, "Reports", "4 generated PDFs + implementation bundle", TEAL)
    d.card(slide, 6.1, 1.0, 2.4, 1.1, "Interactive Demo", "Streamlit app with evidence tables", GOLD)
    d.card(slide, 8.75, 1.0, 2.4, 1.1, "Reproducibility", "Local corpus assets, JSON metrics, plots", SUCCESS)
    d.bullet_box(
        slide,
        [
            "The project is unusually strong as a final presentation because every output mode reinforces the same research story",
            "The app exposes predictions, evidence spans, metrics, dataset views, and literature references",
            "The report pipeline ensures the deck, PDFs, and codebase all stay aligned",
        ],
        0.8,
        2.5,
        5.2,
        2.3,
        17,
    )
    d.picture_contain(slide, CONTACT / "analysis_contact_sheet.png", 6.1, 2.15, 6.15, 4.55)
    d.add_notes(
        slide,
        "This slide is valuable in academic evaluation because it shows end-to-end ownership. "
        "The project is not just a model; it is a complete research package with presentation-grade outputs. "
        "Animation suggestion: cards appear together.",
    )

    slide = d.add_slide("Survey Snapshot: Classical SRL Foundations", "Survey")
    d.text_box(slide, "2002", 0.9, 1.15, 0.8, 0.3, 22, GOLD, True, PP_ALIGN.CENTER)
    d.text_box(slide, "Automatic Labeling of Semantic Roles\nGildea & Jurafsky", 0.55, 1.55, 2.0, 0.8, 15, INK, True, PP_ALIGN.CENTER)
    d.text_box(slide, "2017", 3.5, 1.15, 0.8, 0.3, 22, GOLD, True, PP_ALIGN.CENTER)
    d.text_box(slide, "Deep SRL: What Works and What's Next\nHe et al.", 3.1, 1.55, 2.2, 0.8, 15, INK, True, PP_ALIGN.CENTER)
    d.text_box(slide, "2022", 6.4, 1.15, 0.8, 0.3, 22, GOLD, True, PP_ALIGN.CENTER)
    d.text_box(slide, "PropBank Comes of Age\nBonial et al.", 6.0, 1.55, 2.1, 0.8, 15, INK, True, PP_ALIGN.CENTER)
    d.text_box(slide, "2025", 9.2, 1.15, 0.8, 0.3, 22, GOLD, True, PP_ALIGN.CENTER)
    d.text_box(slide, "LLMs Can Also Do Well\nLi et al.", 8.8, 1.55, 2.1, 0.8, 15, INK, True, PP_ALIGN.CENTER)
    d.text_box(slide, "What earlier systems established", 0.8, 3.0, 3.3, 0.3, 18, NAVY, True)
    d.bullet_box(
        slide,
        [
            "Feature-based SRL showed the importance of syntax and predicate-argument structure",
            "Deep SRL reduced manual feature engineering and improved contextual generalization",
            "Modern PropBank work expanded semantic coverage and resource maturity",
            "Recent LLM-based SRL shows strong promise, but task-specific adaptation still matters",
        ],
        0.9,
        3.45,
        5.1,
        2.4,
        16,
    )
    d.picture_contain(slide, CONTACT / "survey_contact_sheet.png", 6.0, 2.7, 6.2, 3.9)
    d.add_notes(
        slide,
        "This slide uses live-checked literature anchors and the local survey PDF. "
        "The goal is not to teach the full field, but to show the progression from classical SRL to modern semantic modeling. "
        "External sources verified during this build: ACL Anthology, arXiv, and dblp entries. "
        "Animation suggestion: timeline points appear one by one.",
    )

    slide = d.add_slide("Survey Snapshot: QA, QA-SRL, and Transformers", "Survey")
    d.add_table(
        slide,
        [
            ["Research Direction", "What it contributed", "Why it matters here"],
            ["Extractive QA", "Efficient span prediction with EM/F1 evaluation", "Our baseline keeps this extraction framing"],
            ["QA-SRL", "Question-answer supervision over semantic roles", "Our data pipeline operationalizes this bridge with PropBank"],
            ["BERT-style QA", "Powerful contextual span proposals", "The hybrid layer uses this as optional support"],
            ["QA-based semantics", "Semi-structured meaning via QA forms", "Supports the project's role-aware question framing"],
        ],
        0.75,
        1.05,
        7.15,
        3.75,
        11,
    )
    d.bullet_box(
        slide,
        [
            "The project does not reject extractive QA; it enriches it with semantic-role structure",
            "QA-SRL is the clearest conceptual precedent for turning roles into questions",
            "Transformer QA is used selectively in the hybrid layer rather than as the only backbone",
        ],
        0.9,
        5.15,
        6.95,
        1.5,
        16,
    )
    d.picture_contain(slide, CONTACT / "research_paper_contact_sheet.png", 8.1, 1.05, 4.3, 5.9)
    d.add_notes(
        slide,
        "Use this slide to show continuity between earlier QA research and the current project. "
        "The main point is that the project combines extractive QA practice with structured semantic supervision. "
        "Animation suggestion: table first, image second.",
    )

    slide = d.add_slide("Survey Snapshot: Latest Structured-Semantics Trends", "Survey")
    d.add_table(
        slide,
        [
            ["Latest Direction", "Live-checked source", "Key takeaway"],
            ["LLM limits on structured semantics", "arXiv 2024", "LLMs show promise, but structured-role errors remain systematic"],
            ["LLM advances on SRL", "ACL Findings 2025", "Fine-tuned LLMs improve SRL, but adaptation is essential"],
            ["Modern PropBank resource evolution", "*SEM 2022", "Semantic resources remain foundational despite model shifts"],
            ["Large-scale QA-SRL parsing", "ACL 2018", "Question-based semantics remains a key bridge"],
        ],
        0.75,
        1.0,
        7.4,
        4.1,
        11,
    )
    d.bullet_box(
        slide,
        [
            "Recent literature does not eliminate the need for explicit structure; it changes how structure can be modeled",
            "This project aligns with that trend by keeping explicit semantic roles and adding optional modern inference support",
            "The deck therefore positions the system as hybrid and research-aware rather than anti-LLM or purely classical",
        ],
        0.85,
        5.35,
        7.2,
        1.2,
        15,
    )
    d.text_box(slide, "External sources were verified during deck construction via ACL Anthology and arXiv searches on April 7, 2026.", 8.35, 1.15, 4.0, 0.7, 14, SLATE)
    d.picture_contain(slide, CONTACT / "innovation_contact_sheet.png", 8.2, 2.0, 4.2, 4.5)
    d.add_notes(
        slide,
        "This slide is where you explicitly show that the literature review is current. "
        "State that the project acknowledges both the progress and the limitations of LLM-based SRL. "
        "Animation suggestion: minimal fade only.",
    )

    slide = d.add_slide("Existing Systems vs Our System", "Comparison")
    d.add_table(
        slide,
        [
            ["System Family", "Strength", "Limitation", "Our Position"],
            ["Classical SRL", "Interpretability, syntax awareness", "Feature engineering and pipeline brittleness", "Keep semantic transparency, learn more jointly"],
            ["Neural SRL", "Better learned context representations", "Still role-focused rather than QA-facing", "Use multi-task learning with answer extraction"],
            ["Extractive QA", "Strong span prediction and mature metrics", "No explicit answer-role interpretation", "Keep span decoding, add role grounding"],
            ["QA-SRL", "Natural question-answer view of semantics", "Not this exact integrated local pipeline", "Operationalize the bridge with PropBank-derived supervision"],
            ["Transformer QA", "Strong contextual answer proposals", "Higher runtime and weaker transparency", "Use as optional support in a hybrid layer"],
            ["LLM-based SRL", "Strong general language competence", "Needs adaptation and can remain structurally inconsistent", "Acknowledge trend, keep deterministic structured core"],
        ],
        0.65,
        1.05,
        12.0,
        4.6,
        11,
    )
    d.bullet_box(
        slide,
        [
            "The project is best understood as an integrated SRL-QA research package rather than as a single-task model",
            "It borrows strengths from prior systems but centers semantic explainability and reproducibility",
        ],
        0.8,
        5.9,
        11.7,
        0.7,
        16,
    )
    d.add_notes(
        slide,
        "This slide is extremely useful in viva-style questioning. "
        "It shows that the project knows where it sits in the literature and what trade-offs it intentionally makes. "
        "Animation suggestion: none; the table should stay readable.",
    )

    slide = d.add_slide("What Our System Adds Beyond Earlier Work", "Comparison")
    d.card(slide, 0.8, 1.0, 2.7, 1.2, "Semantic grounding", "Every answer is aligned with predicate-argument structure", TEAL)
    d.card(slide, 3.75, 1.0, 2.7, 1.2, "Hybrid reasoning", "Role-aware reranking improves targeted question types", BLUE)
    d.card(slide, 6.7, 1.0, 2.7, 1.2, "Research packaging", "Benchmarks, plots, PDFs, and app all stay aligned", GOLD)
    d.card(slide, 9.65, 1.0, 2.7, 1.2, "Offline reproducibility", "Local assets and saved outputs support reruns", SUCCESS)
    d.bullet_box(
        slide,
        [
            "Earlier systems usually optimized one axis strongly: either semantic structure, or answer extraction, or usability",
            "This project aims for a balanced combination of semantic fidelity, extractive QA, interpretability, and presentation readiness",
            "That balance is what makes it suitable for a final-year research presentation",
        ],
        0.85,
        2.7,
        5.35,
        2.0,
        17,
    )
    d.picture_contain(slide, PLOTS / "research_architecture.png", 6.2, 2.45, 6.05, 4.1)
    d.add_notes(
        slide,
        "Use this slide to translate the comparison table into a simpler story. "
        'The phrase to emphasize is "balanced combination" rather than "state-of-the-art replacement." '
        "Animation suggestion: cards appear together, then architecture.",
    )

    slide = d.add_slide("Strengths, Limitations, and Validity Threats", "Analysis")
    d.text_box(slide, "Strengths", 0.85, 1.0, 2.0, 0.3, 18, SUCCESS, True)
    d.bullet_box(
        slide,
        [
            "Real PropBank-backed semantic supervision",
            "Transparent and modular codebase",
            "Strong artifact trail: metrics, plots, PDFs, app",
            "Targeted hybrid gains on role-sensitive questions",
        ],
        0.85,
        1.4,
        3.8,
        2.3,
        16,
    )
    d.text_box(slide, "Limitations", 4.75, 1.0, 2.0, 0.3, 18, WARN, True)
    d.bullet_box(
        slide,
        [
            "Long-tail PropBank labels remain weak",
            "Question generation is template-based",
            "Hybrid gains are not uniform across every category",
            "Transformer support is expensive on local CPU inference",
        ],
        4.75,
        1.4,
        3.8,
        2.3,
        16,
    )
    d.text_box(slide, "Threats to validity", 8.65, 1.0, 2.5, 0.3, 18, BLUE, True)
    d.bullet_box(
        slide,
        [
            "Treebank-backed subset may differ from the full PropBank distribution",
            "Curated challenge suite is small by design",
            "Some benchmark gains may reflect targeted role heuristics more than global robustness",
        ],
        8.65,
        1.4,
        3.7,
        2.3,
        16,
    )
    d.picture_contain(slide, PLOTS / "confidence_histogram.png", 1.0, 4.2, 4.3, 2.2)
    d.picture_contain(slide, PLOTS / "frame_graph.png", 5.0, 4.2, 3.5, 2.2)
    d.picture_contain(slide, PLOTS / "frame_memory.png", 8.9, 4.2, 3.0, 2.2)
    d.add_notes(
        slide,
        "This slide adds credibility by being explicit about where the system is strong and where it still needs work. "
        "You can keep the discussion short during the talk, but it is useful to have in the deck for questioning. "
        "Animation suggestion: no animation; this is a reference slide.",
    )

    slide = d.add_slide("Future Work and Research Roadmap", "Future")
    d.add_table(
        slide,
        [
            ["Roadmap Step", "Motivation", "Likely Impact"],
            ["Expand beyond local Treebank-backed subset", "Increase coverage and diversity", "Better generalization and stronger evidence"],
            ["Improve paraphrastic question generation", "Reduce template rigidity", "Better QA robustness"],
            ["Target rare-role modeling", "Lift macro F1 and long-tail behavior", "More balanced SRL performance"],
            ["Add causal and temporal reasoning improvements", "Current WHY and some WHEN cases remain hard", "Better semantic completeness"],
            ["Make transformer use selective", "Reduce latency overhead", "More practical hybrid deployment"],
            ["Add human evaluation for explanations", "Validate reasoning summaries and usability", "Stronger presentation and research claims"],
        ],
        0.65,
        1.05,
        12.0,
        4.5,
        10,
    )
    d.text_box(slide, "The current repository already contains the modular pieces needed for most of these extensions.", 0.9, 5.9, 11.0, 0.45, 16, NAVY, True)
    d.add_notes(
        slide,
        "Present future work as a continuation of the same research program rather than as a list of unrelated ideas. "
        "Each next step should feel like a natural extension of the existing modular repo. "
        "Animation suggestion: table only.",
    )

    slide = d.add_slide("PDF Deliverable Gallery: Survey Report", "Appendix")
    d.bullet_box(
        slide,
        [
            "The survey PDF consolidates the literature framing used in the later survey slides",
            "Rendering the report pages inside the deck shows how the repository already packages the project in academic format",
            "This gallery also ensures the final PPT visibly reuses the generated PDF artifacts",
        ],
        0.75,
        1.0,
        4.9,
        2.3,
        16,
    )
    d.picture_contain(slide, CONTACT / "survey_contact_sheet.png", 5.55, 0.95, 6.7, 5.95)
    d.add_notes(
        slide,
        "This is an appendix-style validation slide. "
        "Use it only if you want to show that the survey report exists and is visually integrated into the project package. "
        "Animation suggestion: none.",
    )

    slide = d.add_slide("PDF Deliverable Gallery: Analysis Report", "Appendix")
    d.bullet_box(
        slide,
        [
            "The analysis PDF captures results, diagnostics, plots, and failure patterns from the local run",
            "Its pages reuse the same metrics and figures shown earlier in the presentation",
            "This alignment between PDF and PPT is part of the repo's reproducibility story",
        ],
        0.75,
        1.0,
        4.9,
        2.3,
        16,
    )
    d.picture_contain(slide, CONTACT / "analysis_contact_sheet.png", 5.55, 0.95, 6.7, 5.95)
    d.add_notes(
        slide,
        "Use this slide as proof that the analysis is not hand-waved; it exists as a generated report and matches the deck. "
        "Animation suggestion: none.",
    )

    slide = d.add_slide("PDF Deliverable Gallery: Innovation and Paper", "Appendix")
    d.bullet_box(
        slide,
        [
            "The innovation report focuses on system-level novelty and design rationale",
            "The short research paper condenses the project into a paper-style academic format",
            "Together, these artifacts complete the submission-ready research package",
        ],
        0.75,
        1.0,
        4.9,
        2.3,
        16,
    )
    d.picture_contain(slide, CONTACT / "innovation_contact_sheet.png", 5.55, 0.95, 3.25, 5.95)
    d.picture_contain(slide, CONTACT / "research_paper_contact_sheet.png", 8.95, 0.95, 3.25, 5.95)
    d.add_notes(
        slide,
        "This slide closes the artifact gallery by showing the remaining generated reports. "
        "If time is short in the actual defense, keep this slide as backup material. "
        "Animation suggestion: none.",
    )

    slide = d.add_slide("Final Takeaways and References", "Close")
    d.text_box(slide, "Key takeaways", 0.8, 1.0, 2.5, 0.3, 18, NAVY, True)
    d.bullet_box(
        slide,
        [
            "PropQA-Net makes extractive answers more informative by grounding them in semantic roles",
            "The hybrid layer improves targeted role-sensitive questions while keeping the decision path explainable",
            "The strongest practical upgrade is the lightweight heuristic-semantic reranker",
            "The repo is ready for academic evaluation because code, evidence, and presentation artifacts are aligned",
        ],
        0.85,
        1.35,
        5.5,
        2.5,
        16,
    )
    d.text_box(slide, "External references used in survey/comparison slides", 0.8, 4.15, 4.6, 0.3, 16, NAVY, True)
    d.bullet_box(
        slide,
        [
            "Gildea & Jurafsky (2002) — Automatic Labeling of Semantic Roles",
            "FitzGerald et al. (2018) — Large-Scale QA-SRL Parsing",
            "Bonial et al. (2022) — PropBank Comes of Age",
            "Cheng et al. (2024) — Potential and Limitations of LLMs in Capturing Structured Semantics",
            "Li et al. (2025) — LLMs Can Also Do Well! Breaking Barriers in SRL",
            "Devlin et al. (2018/2019) — BERT",
        ],
        0.85,
        4.5,
        6.3,
        2.0,
        13,
    )
    d.add_table(slide, [["Title", "URL"]] + [[title, url] for title, url in ext_sources[:4]], 7.35, 1.1, 5.2, 5.2, 10)
    d.add_notes(
        slide,
        "Close by reinforcing the core contribution: semantic-role-aware answers are more useful than raw span extraction alone. "
        "Mention that external literature sources were live-checked during deck construction and the rest of the evidence comes from the local repository. "
        "Animation suggestion: keep the closing static and professional.",
    )

    d.save(PPT_PATH)
    print(f"Created: {PPT_PATH}")
    print(f"Slides: {len(d.prs.slides)}")


if __name__ == "__main__":
    build_deck()
